"""
Screenshots 模式：把每张幻灯片截图嵌成 PPT 里的全幅 image。

对应 gen-pptx 的 src/render/build-screenshot.ts。

实现要点：
- python-pptx 没有 PptxGenJS 的 defineLayout；用 prs.slide_width / slide_height 直接设
  Inches（自定义 layout 不需要预先 defineLayout，python-pptx 会按 deck 实际尺寸渲染）
- 每张 slide 嵌一张图片：x=0 y=0 w=slide_w h=slide_h
- speaker notes：slide.notes_slide.notes_text_frame.text = ...
"""

from __future__ import annotations

import base64
import io
from typing import List

from pptx import Presentation
from pptx.util import Inches

# R17-3: 之前从 agentmatrix.desktop.pptx_utils.ooxml 导入 px_to_inches，
# 但 pptx_export 内部其他 render 模块（editable.py 等）统一用 ..parse.units.px_to_inches。
# 两个实现都是 DPI=96 等价，但导入路径分叉会让后续修改时容易漏同步。
from ..parse.units import px_to_inches


def build_screenshot_pptx(image_data_urls: List[str],
                          notes: List[str],
                          width_px: int,
                          height_px: int) -> dict:
    """从一组全屏截图 data URL 生成 .pptx bytes。

    Args:
        image_data_urls: 每张幻灯片的 PNG data URL（"data:image/png;base64,..."）
        notes: 每张幻灯片的讲稿，按 index 对齐，越界 / 空串都视为无讲稿
        width_px / height_px: 幻灯片尺寸（px），用来定义 deck layout

    Returns:
        dict(bytes=b..., bytes_len=int, slides=int, warnings=[...])
    """
    warnings: List[str] = []
    prs = Presentation()
    # 16:9 / 4:3 默认是 EMU；这里直接覆盖成我们要的自定义尺寸（Inches）
    prs.slide_width = Inches(px_to_inches(width_px))
    prs.slide_height = Inches(px_to_inches(height_px))

    # 用 blank layout（避免 placeholder 干扰），index 6 是 blank in default template
    blank_layout = prs.slide_layouts[6]

    for i, data_url in enumerate(image_data_urls):
        slide = prs.slides.add_slide(blank_layout)
        png_bytes = _data_url_to_png(data_url)
        if png_bytes is None:
            # 与 TS 行为对齐：TS build-screenshot.ts 不做 data URL 验证，永远
            # addSlide() 保持页数一致；坏 data 交给 PptxGenJS 自己处理。Python
            # 这里 add_slide 已经先执行（保证页数 == image_data_urls 长度），
            # 解析失败时仅留空白页 + 警告，但仍照常写 speaker notes。
            warnings.append(f"slide {i + 1}: data URL 解析失败，保留空白页")
        else:
            slide.shapes.add_picture(
                io.BytesIO(png_bytes),
                left=0, top=0,
                width=prs.slide_width, height=prs.slide_height,
            )
        # R17-5: 同 R16-7 —— 仅用 strip() 判空，写入保留原文（与 TS 对齐）。
        raw_note = notes[i] if (i < len(notes) and notes[i]) else ""
        if raw_note.strip():
            try:
                slide.notes_slide.notes_text_frame.text = raw_note
            except Exception as e:
                warnings.append(f"addNotes slide {i + 1}: {e}")

    buf = io.BytesIO()
    prs.save(buf)
    raw = buf.getvalue()
    return {
        "bytes": raw,
        "bytes_len": len(raw),
        "slides": len(image_data_urls),
        "warnings": warnings,
    }


def _data_url_to_png(data_url: str) -> bytes:
    if not data_url or not data_url.startswith("data:"):
        return None
    try:
        # data:image/png;base64,XXXX
        header, b64 = data_url.split(",", 1)
        return base64.b64decode(b64)
    except Exception:
        return None
