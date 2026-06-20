"""PptxGenJS-style slide API backed by python-pptx. 对应 render/context.ts 的 PptxSlide.

PptxGenJS 和 python-pptx 的 API 差很大，render-node 直接调 python-pptx 会让逻辑
分散难对齐。这里做一个 adapter：保留 .addText / .addShape / .addImage / .addNotes /
.background 形式，opts 字典也按 PptxGenJS 的 key 命名，内部翻译成 python-pptx 调用。

尚未完全支持的能力（Phase 3 会用 lxml 直写 OOXML 补齐）：
- 文本阴影（PptxGenJS shadow）—— 暂记 warning 跳过
- 图片透明度 —— 暂跳过
- underline 颜色 —— 暂用文本颜色近似
- line spacing as multiple 是支持的；lineSpacingMultiple 直接给 paragraph.line_spacing
- bullet / charSpacing / transparency 都做了 lxml 直写
"""

from __future__ import annotations

import base64
import io
from typing import Any, Dict, List, Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu

from agentmatrix.desktop.pptx_utils.ooxml import qn, ratio_to_ooxml_pct
from agentmatrix.desktop.pptx_utils.shapes import get_mso_shape


# ---------------------------------------------------------------------------
# option helpers
# ---------------------------------------------------------------------------

_ALIGN_MAP = {
    'left': PP_ALIGN.LEFT,
    'center': PP_ALIGN.CENTER,
    'right': PP_ALIGN.RIGHT,
    'justify': PP_ALIGN.JUSTIFY,
}

_VALIGN_MAP = {
    'top': MSO_ANCHOR.TOP,
    'middle': MSO_ANCHOR.MIDDLE,
    'bottom': MSO_ANCHOR.BOTTOM,
}


def _hex_to_rgbcolor(hex_str: str) -> RGBColor:
    h = (hex_str or '000000').lstrip('#')[:6].upper()
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _opt_in(opts: Dict, *keys, default=None):
    for k in keys:
        v = opts.get(k)
        if v is not None:
            return v
    return default


def _apply_fill(shape, fill_spec: Optional[Dict], transparency: int = 0):
    """fill: {color, transparency} | {type:'none'} | None."""
    if not fill_spec or fill_spec.get('type') == 'none':
        shape.fill.background()  # no fill
        return
    color_hex = fill_spec.get('color')
    if not color_hex:
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgbcolor(color_hex)
    # transparency (0-100, PptxGenJS style) → alpha (0-1)
    # 注意边界：transparency=100 表示完全透明（alpha=0），必须走注入路径，
    # 不能用严格 < 100 跳过。_inject_color_alpha 内部已 clamp 到 [0,1]。
    fill_trans = fill_spec.get('transparency', transparency)
    if fill_trans and 0 < int(fill_trans) <= 100:
        _inject_color_alpha(shape.fill.fore_color, 1.0 - int(fill_trans) / 100.0)


def _apply_line(shape, line_spec: Optional[Dict]):
    if not line_spec or line_spec.get('type') == 'none':
        shape.line.fill.background()  # no line
        return
    color_hex = line_spec.get('color')
    width_pt = line_spec.get('width')
    if not color_hex:
        shape.line.fill.background()
        return
    shape.line.color.rgb = _hex_to_rgbcolor(color_hex)
    if width_pt:
        shape.line.width = Pt(float(width_pt))
    dash = line_spec.get('dashType')
    if dash:
        _set_line_dash(shape, dash)
    trans = line_spec.get('transparency')
    if trans and 0 < int(trans) <= 100:
        _inject_color_alpha(shape.line.color, 1.0 - int(trans) / 100.0)


def _set_line_dash(shape, dash_type: str):
    """Set prstDash on shape line via lxml."""
    from lxml import etree
    ln = shape.line._get_or_add_ln()
    # remove existing prstDash
    for existing in ln.findall(qn('a:prstDash')):
        ln.remove(existing)
    prst = etree.SubElement(ln, qn('a:prstDash'))
    prst.set('val', 'solid' if dash_type == 'solid' else dash_type)


def _inject_color_alpha(color_format, alpha_0_to_1: float):
    """Inject a:alpha into a python-pptx ColorFormat (srgbClr must already exist)."""
    from lxml import etree
    xfill = color_format._xFill
    srgb = xfill.find(qn('a:srgbClr'))
    if srgb is None:
        return
    existing = srgb.find(qn('a:alpha'))
    if existing is not None:
        srgb.remove(existing)
    alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
    alpha_elem.set('val', str(ratio_to_ooxml_pct(max(0.0, min(1.0, alpha_0_to_1)))))


def _set_char_spacing(run, points: Optional[float]):
    if points is None:
        return
    from ..parse.units import js_round
    # OOXML spc is in 1/100 pt; positive or negative allowed
    run.font._rPr.set('spc', str(js_round(points * 100)))


def _set_bullet(paragraph, bullet_spec):
    """bullet: False | {type:'number', indent, numberStartAt?, numberType?} | {indent, characterCode?} | True."""
    from lxml import etree
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    # remove existing bullet defs
    for tag in ('a:buChar', 'a:buAutoNum', 'a:buNone', 'a:buFont', 'a:buClr', 'a:buSzPct'):
        for existing in pPr.findall(qn(tag)):
            pPr.remove(existing)
    if bullet_spec is False or bullet_spec is None:
        bu_none = etree.SubElement(pPr, qn('a:buNone'))
        return
    # indent (in pt) → buSzPct & pPr marL/indent
    # PptxGenJS 默认 27pt（.d.ts:1118 DEF_BULLET_MARGIN），移植时漏掉了。
    indent_pt = float(bullet_spec.get('indent') or 27)
    indent_emu = int(indent_pt * 12700)
    marL = indent_emu
    pPr.set('marL', str(marL))
    pPr.set('indent', str(-indent_emu))
    bu_font = etree.SubElement(pPr, qn('a:buFont'))
    bu_font.set('typeface', 'Arial')
    if bullet_spec.get('type') == 'number':
        bu = etree.SubElement(pPr, qn('a:buAutoNum'))
        ntype = bullet_spec.get('numberType', 'arabicPeriod')
        # map PptxGenJS numberType → OOXML type
        bu.set('type', _PPTX_NUM_TYPE_MAP.get(ntype, 'arabicPeriod'))
        start = bullet_spec.get('numberStartAt')
        if start:
            bu.set('startAt', str(int(start)))
    else:
        bu = etree.SubElement(pPr, qn('a:buChar'))
        bu.set('char', bullet_spec.get('characterCode') or '•')


_PPTX_NUM_TYPE_MAP = {
    'alphaLcPeriod': 'alphaLcPeriod',
    'alphaUcPeriod': 'alphaUcPeriod',
    'romanLcPeriod': 'romanLcPeriod',
    'romanUcPeriod': 'romanUcPeriod',
    'alphaLcParenBoth': 'alphaLcParenBoth',
    'alphaUcParenBoth': 'alphaUcParenBoth',
    'alphaLcParenR': 'alphaLcParenR',
    'alphaUcParenR': 'alphaUcParenR',
    'romanLcParenBoth': 'romanLcParenBoth',
    'romanUcParenR': 'romanLcParenR',
    'arabicPeriod': 'arabicPeriod',
    'arabicParenR': 'arabicParenR',
    'arabicDbPlain': 'arabicDbPlain',
}


def _set_tab_stops(paragraph, tab_stops: Optional[List[Dict]]):
    if not tab_stops:
        return
    from lxml import etree
    pPr = paragraph._p.get_or_add_pPr()
    for existing in pPr.findall(qn('a:tabLst')):
        pPr.remove(existing)
    if not tab_stops:
        return
    tab_lst = etree.SubElement(pPr, qn('a:tabLst'))
    for tab in tab_stops:
        tab_elem = etree.SubElement(tab_lst, qn('a:tab'))
        position_in = tab.get('position', 0)
        tab_elem.set('pos', str(int(float(position_in) * 914400)))


def _apply_shape_shadow(shape, shadow_spec: Optional[Dict]):
    """shadow spec → a:effectLst/a:outerShdw in shape's spPr.

    Shadow spec fields (PptxGenJS convention):
      opacity 0..1, offset (pt, magnitude), blur (pt), angle (deg, 0=right 90=down),
      color (hex).
    OOXML outerShdw encodes direction as `dir` (1/60000 deg) + scalar `dist`,
    so we decompose the offset magnitude via angle into x/y for the helper
    (which then derives dir/dist itself — no overrides needed).
    """
    if not shadow_spec:
        return
    import math
    from agentmatrix.desktop.pptx_utils.ooxml import (
        make_outer_shadow_xml, make_effect_lst,
    )
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    # remove existing effectLst
    for existing in spPr.findall(qn('a:effectLst')):
        spPr.remove(existing)

    opacity = float(shadow_spec.get('opacity', 1.0))
    offset_pt = float(shadow_spec.get('offset', 0))
    blur_pt = float(shadow_spec.get('blur', 0))
    angle = float(shadow_spec.get('angle', 90))
    color_hex = shadow_spec.get('color', '000000')

    # offset in pt → px (px = pt / 0.75), decomposed by angle for the helper
    offset_px = offset_pt / 0.75
    blur_px = blur_pt / 0.75
    rad = math.radians(angle % 360)
    offset_x_px = offset_px * math.cos(rad)
    offset_y_px = offset_px * math.sin(rad)

    shadow_xml = make_outer_shadow_xml(
        blur_px=blur_px,
        offset_x_px=offset_x_px,
        offset_y_px=offset_y_px,
        color_hex=color_hex,
        alpha=opacity,
    )
    spPr.append(make_effect_lst(shadow_xml))


# ---------------------------------------------------------------------------
# Adapter
# ---------------------------------------------------------------------------

class PptxSlideAdapter:
    """PptxGenJS-style API backed by a python-pptx slide.

    Wrap with prs.slide_layouts[6] (blank). Coordinates are in inches.
    """

    def __init__(self, prs, blank_layout=None):
        if blank_layout is None:
            blank_layout = prs.slide_layouts[6]
        self._prs = prs
        self._slide = prs.slides.add_slide(blank_layout)
        self._shapes = self._slide.shapes

    # ---------- background ----------
    @property
    def background(self):
        return None

    @background.setter
    def background(self, value: Optional[Dict]):
        if not value:
            return
        color = value.get('color')
        if not color:
            return
        self._slide.background.fill.solid()
        self._slide.background.fill.fore_color.rgb = _hex_to_rgbcolor(color)

    # ---------- text ----------
    def addText(self, text, opts: Dict):
        """Add a text box.

        text: str OR list of {text: str, options: {...run-level...}}.
        opts: shape-level options (x/y/w/h/margin/fontSize/align/valign/wrap/fit/
              lineSpacingMultiple/charSpacing/rotate/shadow/hyperlink/bullet/
              tabStops/fontFace/bold/italic/...).
        """
        x = Inches(float(opts.get('x', 0)))
        y = Inches(float(opts.get('y', 0)))
        w = Inches(float(opts.get('w', 0)))
        h = Inches(float(opts.get('h', 0)))
        tb = self._shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        margin = opts.get('margin')
        if margin is not None:
            # margin 单位是 pt（与 PptxGenJS 一致），不是 inch；
            # 1 pt = 12700 EMU。之前误用 914400 (EMU/inch) 导致 2pt 被当成 2 inch，
            # 每边吃掉 192px，textbbox 实际可用宽度骤减 → 文字疯狂换行。
            # PptxGenJS margin 支持 number 或 [top,right,bottom,left]（.d.ts:1851）。
            if isinstance(margin, (list, tuple)) and len(margin) == 4:
                m_top, m_right, m_bottom, m_left = margin
            else:
                m_top = m_right = m_bottom = m_left = margin
            tf.margin_left = Emu(int(float(m_left) * 12700))
            tf.margin_right = Emu(int(float(m_right) * 12700))
            tf.margin_top = Emu(int(float(m_top) * 12700))
            tf.margin_bottom = Emu(int(float(m_bottom) * 12700))

        wrap = opts.get('wrap', True)
        tf.word_wrap = bool(wrap)
        if opts.get('fit') == 'shrink':
            try:
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass

        # valign on textbox
        # PptxGenJS TextPropsOptions.valign 默认 'middle'（.d.ts:1883 @default middle）。
        # 上层 render-node 目前始终显式传 valign，但 API 契约缺失时 fallback
        # 必须与 PptxGenJS 对齐，否则裸调用方会得到 top（HTML/CSS 直觉）。
        valign = opts.get('valign', 'middle')
        try:
            tf.vertical_anchor = _VALIGN_MAP.get(valign, MSO_ANCHOR.MIDDLE)
        except Exception:
            pass

        if opts.get('rotate') is not None:
            try:
                tb.rotation = float(opts['rotate'])
            except Exception:
                pass

        # Normalize to list of (text, run_opts) tuples.
        if isinstance(text, str):
            text_objs = [(text, {})]
        else:
            text_objs = [(obj.get('text', ''), obj.get('options') or {}) for obj in text]

        # First paragraph already exists (empty). Use it for the first run.
        first_para = tf.paragraphs[0]
        # PptxGenJS semantics: breakLine=true on a run means "this run ends with
        # a line break" — i.e. the NEXT run starts a new paragraph. Default is
        # False (so multi-run text without explicit breakLine stays on one line).
        para = first_para
        _apply_paragraph_opts(first_para, opts, text_objs[0][1] if text_objs else None)
        prev_break_line = False
        for i, (txt, run_opts) in enumerate(text_objs):
            if i > 0 and prev_break_line:
                para = tf.add_paragraph()
                _apply_paragraph_opts(para, opts, run_opts)
            # Merge run-level + shape-level font settings
            merged_opts = {**opts, **run_opts}
            run = para.add_run()
            run.text = txt
            _apply_run_font(run, merged_opts)
            # Hyperlink: shape-level (opts) is the default; per-run overrides.
            href = merged_opts.get('hyperlink')
            if href and isinstance(href, dict):
                url = href.get('url')
                if url:
                    run.hyperlink.address = url
            prev_break_line = bool(run_opts.get('breakLine', False))

        return tb

    # ---------- shape ----------
    def addShape(self, shape_name: str, opts: Dict):
        mso = get_mso_shape(shape_name, fallback='rect')
        if mso is None:
            # fallback to RECTANGLE
            from pptx.enum.shapes import MSO_SHAPE as _MS
            mso = _MS.RECTANGLE
        x = Inches(float(opts.get('x', 0)))
        y = Inches(float(opts.get('y', 0)))
        w = Inches(float(opts.get('w', 0)))
        h = Inches(float(opts.get('h', 0)))
        shape = self._shapes.add_shape(mso, x, y, w, h)

        _apply_fill(shape, opts.get('fill'))
        _apply_line(shape, opts.get('line'))
        _apply_shape_shadow(shape, opts.get('shadow'))

        # rounded rect radius
        radius = opts.get('rectRadius')
        if radius and shape_name in ('roundRect', 'rounded-rectangle'):
            try:
                min_side_in = min(float(opts.get('w', 0)), float(opts.get('h', 0)))
                if min_side_in > 0:
                    ratio = max(0.0, min(0.5, float(radius) / min_side_in))
                    shape.adjustments[0] = ratio
            except Exception:
                pass

        if opts.get('rotate') is not None:
            try:
                shape.rotation = float(opts['rotate'])
            except Exception:
                pass

        # remove the default text content python-pptx adds
        try:
            shape.text_frame.text = ''
        except Exception:
            pass
        return shape

    # ---------- image ----------
    def addImage(self, opts: Dict):
        x = Inches(float(opts.get('x', 0)))
        y = Inches(float(opts.get('y', 0)))
        w = Inches(float(opts.get('w', 0)))
        h = Inches(float(opts.get('h', 0)))

        # PptxGenJS 支持 path 或 data（DataOrPathProps）。本移植目前只支持 data
        # URL —— 上层 render-node 的 resolve_media 已把所有 URL/SVG 转 data URL，
        # 实际调用方从不传 path。若未来 caller 直接传 path，需要在这里加：
        #   path_url = opts.get('path')
        #   if path_url and not data_url: <fetch/local-read → data URL>
        # 见 finding A6-1（已降级为 nitpick，因当前所有 caller 均传 data）。
        data_url = opts.get('data') or opts.get('dataUrl') or ''
        if not data_url or not data_url.startswith('data:'):
            # A6-2: 之前静默 return None，导致图莫名其妙消失且无 warning。
            # 改为 raise —— caller (node.py:252) 有 try/except，会转成 warning。
            raise ValueError(
                f"addImage: data URL 缺失或非 data: 协议（got {data_url[:40]!r}...）"
            )
        try:
            header, b64 = data_url.split(',', 1)
            raw = base64.b64decode(b64)
        except Exception as e:
            raise ValueError(f"addImage: data URL base64 解码失败：{e}") from e
        buf = io.BytesIO(raw)

        # default sizing = stretch (python-pptx default for add_picture with w/h)
        sizing = opts.get('sizing')
        sizing_type = sizing.get('type') if sizing else None
        if sizing_type in ('cover', 'contain', 'crop'):
            # A6-9: PptxGenJS sizing.w/h 是必填（.d.ts:1356/1362），缺失应报错
            # 而非静默 fallback 到 stretch（输出尺寸会与期望不符）。
            if not sizing.get('w') or not sizing.get('h'):
                raise ValueError(
                    f"addImage: sizing.type={sizing_type!r} 要求 w/h，"
                    f"got w={sizing.get('w')!r} h={sizing.get('h')!r}"
                )
            if sizing_type == 'crop':
                # PptxGenJS sizing.crop（.d.ts:1349）：把 blip 按 sizing.{x,y,w,h}
                # 裁出子矩形，填回原 box（不缩放）。用 a:srcRect 实现。
                pic = self._shapes.add_picture(buf, x, y, width=w, height=h)
                _apply_picture_crop_rect(pic, sizing)
            else:
                # cover/contain：先按 sizing.w/h 还原 intrinsic 尺寸，再 fillRect 裁/留白
                pic = self._shapes.add_picture(
                    buf, Inches(float(opts.get('x', 0))), Inches(float(opts.get('y', 0))),
                    width=Inches(float(sizing['w'])), height=Inches(float(sizing['h'])),
                )
                _apply_picture_cover_contain(pic, sizing_type, float(opts.get('w', 0)), float(opts.get('h', 0)))
        else:
            pic = self._shapes.add_picture(buf, x, y, width=w, height=h)

        if opts.get('rounding'):
            _apply_picture_rounding(pic)

        # flipH / flipV（PptxGenJS ImageProps .d.ts:1308-1314）：在 spPr/xfrm 上
        # 设 'flipH' / 'flipV' 属性（OOXML a:xfrm）。
        flip_h = opts.get('flipH')
        flip_v = opts.get('flipV')
        if flip_h or flip_v:
            try:
                _apply_picture_flip(pic, bool(flip_h), bool(flip_v))
            except Exception:
                pass

        if opts.get('rotate') is not None:
            try:
                pic.rotation = float(opts['rotate'])
            except Exception:
                pass

        # hyperlink（PptxGenJS ImageProps.hyperlink .d.ts:1315）：图片可点击跳转。
        # python-pptx pic.click_action.target_url 即可。
        hyperlink = opts.get('hyperlink')
        if hyperlink and isinstance(hyperlink, dict):
            url = hyperlink.get('url')
            if url:
                try:
                    pic.click_action.target_url = url
                except Exception:
                    pass

        trans = opts.get('transparency')
        if trans and 0 < int(trans) <= 100:
            _apply_picture_transparency(pic, 1.0 - int(trans) / 100.0)
        return pic

    # ---------- notes ----------
    def addNotes(self, notes: str):
        if not notes:
            return
        self._slide.notes_slide.notes_text_frame.text = notes


# ---------------------------------------------------------------------------
# Internal: paragraph / run formatting
# ---------------------------------------------------------------------------

def _apply_paragraph_opts(paragraph, opts: Dict, run_opts: Optional[Dict] = None):
    """Apply paragraph-level options. Per-run (run_opts) overrides shape-level for
    bullet / tabStops, since lists put per-item bullets there."""
    merged = dict(opts)
    if run_opts:
        for k in ('bullet', 'tabStops'):
            if k in run_opts:
                merged[k] = run_opts[k]
    align = merged.get('align')
    if align:
        try:
            paragraph.alignment = _ALIGN_MAP.get(align, PP_ALIGN.LEFT)
        except Exception:
            pass
    lsp = merged.get('lineSpacingMultiple')
    if lsp is not None:
        try:
            paragraph.line_spacing = float(lsp)
        except Exception:
            pass
    if 'bullet' in merged:
        try:
            _set_bullet(paragraph, merged['bullet'])
        except Exception:
            pass
    if merged.get('tabStops'):
        try:
            _set_tab_stops(paragraph, merged['tabStops'])
        except Exception:
            pass


def _apply_run_font(run, opts: Dict):
    font = run.font
    ff = opts.get('fontFace')
    if ff:
        font.name = ff
    fs = opts.get('fontSize')
    if fs is not None:
        try:
            font.size = Pt(float(fs))
        except Exception:
            pass
    font.bold = bool(opts.get('bold')) or None
    font.italic = bool(opts.get('italic')) or None
    # PptxGenJS strike: boolean | 'sngStrike' | 'dblStrike'
    # （.d.ts:1878）。true / 'sngStrike' → 单删除线；'dblStrike' → 双删除线。
    s = opts.get('strike')
    if s:
        try:
            run.font._rPr.set('strike', 'dblStrike' if s == 'dblStrike' else 'sngStrike')
        except Exception:
            pass
    if opts.get('underline'):
        u = opts['underline']
        if isinstance(u, dict):
            style = u.get('style', 'sng')
            underline_color = u.get('color')
        else:
            style = 'sng'
            underline_color = None
        from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE
        u_map = {
            'sng': MSO_TEXT_UNDERLINE_TYPE.SINGLE_LINE,
            'dbl': MSO_TEXT_UNDERLINE_TYPE.DOUBLE_LINE,
            'dotted': MSO_TEXT_UNDERLINE_TYPE.DOTTED_LINE,
            'dash': MSO_TEXT_UNDERLINE_TYPE.DASH_LINE,
            'wavy': MSO_TEXT_UNDERLINE_TYPE.WAVY_LINE,
        }
        try:
            font.underline = u_map.get(style, MSO_TEXT_UNDERLINE_TYPE.SINGLE_LINE)
        except Exception:
            pass
        # PptxGenJS underline.color（.d.ts:1254）—— python-pptx 没有公开 API，
        # 直接 lxml 写 rPr 下的 a:uFill/a:srgbClr。
        if underline_color:
            try:
                _set_underline_color(run, underline_color)
            except Exception:
                pass

    color_hex = opts.get('color')
    if color_hex:
        font.color.rgb = _hex_to_rgbcolor(color_hex)
        trans = opts.get('transparency')
        if trans and 0 < int(trans) <= 100:
            _inject_color_alpha(font.color, 1.0 - int(trans) / 100.0)

    if opts.get('subscript'):
        font.subscript = True
    if opts.get('superscript'):
        font.superscript = True

    if opts.get('highlight'):
        try:
            _apply_run_highlight(run, opts['highlight'])
        except Exception:
            pass

    if opts.get('charSpacing') is not None:
        try:
            _set_char_spacing(run, opts['charSpacing'])
        except Exception:
            pass

    shadow = opts.get('shadow')
    if shadow:
        try:
            _apply_run_shadow(run, shadow)
        except Exception:
            pass


def _apply_run_shadow(run, shadow_spec: Dict):
    """Apply text shadow as a:effectLst/a:outerShdw on the run's rPr.

    PptxGenJS applies text shadows at the run level (a:rPr/a:effectLst),
    NOT as a shape-level spPr effect. OOXML supports per-run effects.

    Schema-order: CT_TextCharacterProperties requires effectLst to come
    AFTER the fill group (solidFill/gradFill/...) and BEFORE a:highlight,
    the underline-fill group, and a:latin/a:ea/a:cs. We locate the
    insertion point by scanning for the first child that should come
    after effectLst; if none, append.
    """
    import math
    from lxml import etree
    from agentmatrix.desktop.pptx_utils.ooxml import (
        make_outer_shadow_xml, make_effect_lst, qn,
    )
    rPr = run.font._rPr
    if rPr is None:
        # add_run() always creates rPr in current python-pptx; be defensive
        rPr = run._r.makeelement(qn('a:rPr'), {})
        run._r.insert(0, rPr)
    # remove existing effectLst on this rPr
    for existing in rPr.findall(qn('a:effectLst')):
        rPr.remove(existing)

    opacity = float(shadow_spec.get('opacity', 1.0))
    offset_pt = float(shadow_spec.get('offset', 0))
    blur_pt = float(shadow_spec.get('blur', 0))
    angle = float(shadow_spec.get('angle', 90))
    color_hex = shadow_spec.get('color', '000000')

    # decompose offset (pt) into x/y using the angle (0=right, 90=down)
    rad = math.radians(angle % 360)
    offset_x_pt = offset_pt * math.cos(rad)
    offset_y_pt = offset_pt * math.sin(rad)
    offset_x_px = offset_x_pt / 0.75
    offset_y_px = offset_y_pt / 0.75
    blur_px = blur_pt / 0.75

    shadow_xml = make_outer_shadow_xml(
        blur_px=blur_px,
        offset_x_px=offset_x_px,
        offset_y_px=offset_y_px,
        color_hex=color_hex,
        alpha=opacity,
    )
    effect_lst = make_effect_lst(shadow_xml)

    # find correct schema-position: insert before any child that should
    # come AFTER effectLst (highlight, underline-fill group, latin/ea/cs, ...)
    _insert_rpr_child_schema_ordered(rPr, effect_lst)


# CT_TextCharacterProperties child schema order (ECMA-376 §20.1.2.2.3)
# Used to insert new rPr children at the correct position.
_RPR_CHILD_ORDER = (
    'ln',
    # fill group (mutually exclusive — any one)
    'noFill', 'solidFill', 'gradFill', 'blipFill', 'pattFill', 'grpFill',
    'effectLst', 'effectDag',
    'highlight',
    # underline-line / underline-fill group
    'uLnTx', 'uLn', 'uFillTx', 'uFill',
    'latin', 'ea', 'cs',
    'sym',
    'hlinkClick', 'hlinkMouseOver',
    'rtl',
    'extLst',
)


def _rpr_child_index(local_name: str) -> int:
    """Schema position for a CT_TextCharacterProperties child local-name.

    Returns len(_RPR_CHILD_ORDER) for unknown tags (so they sort last).
    """
    try:
        return _RPR_CHILD_ORDER.index(local_name)
    except ValueError:
        return len(_RPR_CHILD_ORDER)


def _insert_rpr_child_schema_ordered(rPr, new_elem):
    """Insert new_elem into rPr at the schema-correct position.

    Walks existing children; inserts before the first one whose schema
    index is greater than new_elem's. If none, appends at end.
    """
    from lxml import etree
    new_idx = _rpr_child_index(etree.QName(new_elem).localname)
    for child in rPr:
        child_idx = _rpr_child_index(etree.QName(child).localname)
        if child_idx > new_idx:
            child.addprevious(new_elem)
            return
    rPr.append(new_elem)


def _apply_run_highlight(run, color_hex):
    """Inject a:highlight/a:srgbClr into run's rPr (python-pptx has no API).

    Per CT_TextCharacterProperties schema, a:highlight comes AFTER
    a:effectLst and BEFORE the underline-fill group; we delegate to
    _insert_rpr_child_schema_ordered for correct placement.
    """
    from lxml import etree
    from agentmatrix.desktop.pptx_utils.ooxml import qn
    rPr = run.font._rPr
    if rPr is None:
        rPr = run._r.makeelement(qn('a:rPr'), {})
        run._r.insert(0, rPr)
    # remove existing highlight
    for existing in rPr.findall(qn('a:highlight')):
        rPr.remove(existing)
    highlight = etree.Element(qn('a:highlight'))
    srgb = etree.SubElement(highlight, qn('a:srgbClr'))
    srgb.set('val', (color_hex or '000000').lstrip('#')[:6].upper())
    _insert_rpr_child_schema_ordered(rPr, highlight)


def _set_underline_color(run, color_hex):
    """Inject a:uFill/a:srgbClr into run's rPr (python-pptx has no API).

    Per CT_TextCharacterProperties schema, the underline-fill group sits
    between highlight and latin/ea/cs. We use _insert_rpr_child_schema_ordered
    for correct placement (uFill is in the schema order list).
    """
    from lxml import etree
    from agentmatrix.desktop.pptx_utils.ooxml import qn
    rPr = run.font._rPr
    if rPr is None:
        rPr = run._r.makeelement(qn('a:rPr'), {})
        run._r.insert(0, rPr)
    # remove existing uFill
    for existing in rPr.findall(qn('a:uFill')):
        rPr.remove(existing)
    ufill = etree.Element(qn('a:uFill'))
    solid = etree.SubElement(ufill, qn('a:solidFill'))
    srgb = etree.SubElement(solid, qn('a:srgbClr'))
    srgb.set('val', (color_hex or '000000').lstrip('#')[:6].upper())
    _insert_rpr_child_schema_ordered(rPr, ufill)


def _apply_picture_cover_contain(picture, fit: str, box_w_in: float, box_h_in: float):
    """Apply a:blipFill/a:stretch/a:fillRect for cover/contain.

    - cover: blip 溢出 frame，每边裁掉相同比例（负 fillRect insets）
    - contain: blip 缩进 frame，每边留 letterbox（正 fillRect insets）
    """
    from agentmatrix.desktop.pptx_utils.ooxml import (
        set_picture_crop, set_picture_letterbox,
    )
    try:
        native_w = picture.image.size[0] / 96.0  # px → inches (image.size is px)
        native_h = picture.image.size[1] / 96.0
    except Exception:
        return
    if native_w <= 0 or native_h <= 0 or box_w_in <= 0 or box_h_in <= 0:
        return

    if fit == 'cover':
        scale = max(box_w_in / native_w, box_h_in / native_h)
        disp_w = native_w * scale
        disp_h = native_h * scale
        # 裁掉比例：每边 (溢出量 / disp) / 2，对 disp 上算（disp 是 blip 实际显示尺寸）
        crop_w = max(0.0, (disp_w - box_w_in) / disp_w / 2) if disp_w > 0 else 0.0
        crop_h = max(0.0, (disp_h - box_h_in) / disp_h / 2) if disp_h > 0 else 0.0
        if crop_w > 0 or crop_h > 0:
            set_picture_crop(picture, left_ratio=crop_w, right_ratio=crop_w,
                             top_ratio=crop_h, bottom_ratio=crop_h)
    else:  # contain
        scale = min(box_w_in / native_w, box_h_in / native_h)
        disp_w = native_w * scale
        disp_h = native_h * scale
        # letterbox 比例：每边留白 / frame，对 frame 上算（frame 是 spPr/xfrm 的尺寸）
        pad_w = max(0.0, (box_w_in - disp_w) / box_w_in / 2) if box_w_in > 0 else 0.0
        pad_h = max(0.0, (box_h_in - disp_h) / box_h_in / 2) if box_h_in > 0 else 0.0
        if pad_w > 0 or pad_h > 0:
            set_picture_letterbox(picture, left_ratio=pad_w, right_ratio=pad_w,
                                  top_ratio=pad_h, bottom_ratio=pad_h)


def _apply_picture_rounding(picture):
    """Replace picture's geometry with ellipse (border-radius:50% style crop)."""
    from lxml import etree
    from agentmatrix.desktop.pptx_utils.ooxml import qn
    sp = picture._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    # remove existing prstGeom
    for existing in spPr.findall(qn('a:prstGeom')):
        spPr.remove(existing)
    prst = etree.SubElement(spPr, qn('a:prstGeom'))
    prst.set('prst', 'ellipse')
    etree.SubElement(prst, qn('a:avLst'))


def _apply_picture_flip(picture, flip_h: bool, flip_v: bool):
    """Set flipH/flipV on the picture's spPr/xfrm (OOXML a:xfrm attributes)."""
    from agentmatrix.desktop.pptx_utils.ooxml import qn
    sp = picture._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return
    if flip_h:
        xfrm.set('flipH', '1')
    if flip_v:
        xfrm.set('flipV', '1')


def _apply_picture_crop_rect(picture, sizing: Dict):
    """PptxGenJS sizing.crop（type='crop'）：用 a:srcRect 裁出子矩形。

    sizing 形如 `{type:'crop', x, y, w, h}`，单位 inches，相对 native image size：
    - x, y：crop 起点偏移（左 / 上）
    - w, h：crop 出的子图尺寸

    转成 OOXML a:srcRect 的 l/t/r/b（百分比，单位 1/1000），未裁的边为 0。
    若 sizing 没给 x/y/w/h，降级为不裁。
    """
    from lxml import etree
    from agentmatrix.desktop.pptx_utils.ooxml import qn, ratio_to_ooxml_pct
    try:
        native_w_in = picture.image.size[0] / 96.0
        native_h_in = picture.image.size[1] / 96.0
    except Exception:
        return
    if native_w_in <= 0 or native_h_in <= 0:
        return
    x = float(sizing.get('x') or 0)
    y = float(sizing.get('y') or 0)
    w = float(sizing.get('w') or native_w_in)
    h = float(sizing.get('h') or native_h_in)
    # 防御：限制到 native 范围内
    x = max(0.0, min(x, native_w_in))
    y = max(0.0, min(y, native_h_in))
    w = max(0.0, min(w, native_w_in - x))
    h = max(0.0, min(h, native_h_in - y))
    left_ratio = x / native_w_in
    top_ratio = y / native_h_in
    right_ratio = (native_w_in - x - w) / native_w_in
    bottom_ratio = (native_h_in - y - h) / native_h_in

    sp = picture._element
    blip_fill = sp.find(qn('p:blipFill'))
    if blip_fill is None:
        blip_fill = sp.find(qn('a:blipFill'))
    if blip_fill is None:
        return
    # 移除已有 stretch/fillRect/srcRect
    for tag in ('a:stretch', 'a:tile'):
        for existing in blip_fill.findall(qn(tag)):
            blip_fill.remove(existing)
    src_rect = etree.SubElement(blip_fill, qn('a:srcRect'))
    src_rect.set('l', str(ratio_to_ooxml_pct(left_ratio)))
    src_rect.set('t', str(ratio_to_ooxml_pct(top_ratio)))
    src_rect.set('r', str(ratio_to_ooxml_pct(right_ratio)))
    src_rect.set('b', str(ratio_to_ooxml_pct(bottom_ratio)))
    stretch = etree.SubElement(blip_fill, qn('a:stretch'))
    etree.SubElement(stretch, qn('a:fillRect'))


def _apply_picture_transparency(picture, alpha_0_to_1: float):
    """Apply alpha to picture via a:alphaModFix inside a:blipFill/a:blip."""
    from lxml import etree
    from agentmatrix.desktop.pptx_utils.ooxml import qn, ratio_to_ooxml_pct
    sp = picture._element
    blip_fill = sp.find(qn('p:blipFill'))
    if blip_fill is None:
        blip_fill = sp.find(qn('a:blipFill'))
    if blip_fill is None:
        return
    blip = blip_fill.find(qn('a:blip'))
    if blip is None:
        return
    # remove any existing alphaModFix
    for existing in blip.findall(qn('a:alphaModFix')):
        blip.remove(existing)
    mod = etree.SubElement(blip, qn('a:alphaModFix'))
    mod.set('amt', str(ratio_to_ooxml_pct(max(0.0, min(1.0, alpha_0_to_1)))))
