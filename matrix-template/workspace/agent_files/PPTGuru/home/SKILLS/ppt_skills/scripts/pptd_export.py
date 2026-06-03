#!/usr/bin/env python3
"""
PPTD Export — converts .pptd format to .pptx using python-pptx.
Usage: python3 pptd_export.py <input.pptd> [-o output.pptx]
"""

import sys
import os
import re
import math
from pathlib import Path
from urllib.parse import urlparse

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

from pptd_common import (
    px_to_emu, emu_to_px, pt_to_emu, emu_to_pt, load_yaml, load_pptd_to_mso_map,
    find_icon_svg, RESOURCE_DIR, NSMAP, parse_font_family,
    inject_srgb_color, PPTD_DASH_TO_XL, map_dash_style,
    ooxml_angle_to_degrees, degrees_to_ooxml_angle,
    ooxml_pct_to_ratio, ratio_to_ooxml_pct,
)
from pptd_color import hex_to_rgb, resolve_theme_ref, resolve_color
from pptd_text import parse_rich_text, unescape_html


# ---------------------------------------------------------------------------
# XML Factory Functions — build OOXML elements directly, no python-pptx defaults
# ---------------------------------------------------------------------------

def _get_sp_tree(slide):
    """Get the shape tree element for a slide."""
    sp_tree = slide._element.find(qn('p:cSld') + '/' + qn('p:spTree'))
    if sp_tree is None:
        sp_tree = slide._element.find(qn('p:spTree'))
    return sp_tree


def _next_shape_id(slide):
    """Get the next unique shape ID for a slide."""
    sp_tree = _get_sp_tree(slide)
    max_id = 0
    for cNvPr in sp_tree.iter(qn('p:cNvPr')):
        try:
            max_id = max(max_id, int(cNvPr.get('id', 0)))
        except (ValueError, TypeError):
            pass
    return max_id + 1


def _make_sp(x_emu, y_emu, w_emu, h_emu, shape_name='rect', txBox=True, name=None, shape_id=None):
    """Create a <p:sp> element and return it as a python-pptx CT_Shape.

    For textboxes (txBox=True): no style, no fill, bodyPr empty (inherits defaults).
    For shapes (txBox=False): no style, no fill (caller adds fill as needed).

    This avoids python-pptx's unwanted defaults:
    - No <p:style> element (which causes fill from theme)
    - No anchor="ctr" on bodyPr
    - No algn="ctr" on pPr
    """
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    sp = etree.Element(qn('p:sp'), nsmap=nsmap)

    # nvSpPr
    nvSpPr = etree.SubElement(sp, qn('p:nvSpPr'))
    cNvPr = etree.SubElement(nvSpPr, qn('p:cNvPr'))
    cNvPr.set('id', str(shape_id or 0))
    cNvPr.set('name', name or 'Shape')
    cNvSpPr = etree.SubElement(nvSpPr, qn('p:cNvSpPr'))
    if txBox:
        cNvSpPr.set('txBox', '1')
    etree.SubElement(nvSpPr, qn('p:nvPr'))

    # spPr
    spPr = etree.SubElement(sp, qn('p:spPr'))
    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(x_emu))
    off.set('y', str(y_emu))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(w_emu))
    ext.set('cy', str(h_emu))
    prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
    prstGeom.set('prst', shape_name)
    etree.SubElement(prstGeom, qn('a:avLst'))

    # txBody — empty bodyPr (no anchor, no wrap override)
    txBody = etree.SubElement(sp, qn('p:txBody'))
    etree.SubElement(txBody, qn('a:bodyPr'))
    etree.SubElement(txBody, qn('a:lstStyle'))
    etree.SubElement(txBody, qn('a:p'))

    # Convert to python-pptx CT_Shape so it works with slide.shapes
    xml_str = etree.tostring(sp, encoding='unicode')
    return parse_xml(xml_str)


def _make_cxnSp(x1_emu, y1_emu, x2_emu, y2_emu, shape_id=None, name=None):
    """Create a <p:cxnSp> (connector) element."""
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    cxnSp = etree.Element(qn('p:cxnSp'), nsmap=nsmap)

    nvCxnSpPr = etree.SubElement(cxnSp, qn('p:nvCxnSpPr'))
    cNvPr = etree.SubElement(nvCxnSpPr, qn('p:cNvPr'))
    cNvPr.set('id', str(shape_id or 0))
    cNvPr.set('name', name or 'Connector')
    cNvCxnSpPr = etree.SubElement(nvCxnSpPr, qn('p:cNvCxnSpPr'))
    etree.SubElement(nvCxnSpPr, qn('p:nvPr'))

    spPr = etree.SubElement(cxnSp, qn('p:spPr'))
    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    # Connector position: origin at (x1, y1), extent to (x2, y2)
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(min(x1_emu, x2_emu)))
    off.set('y', str(min(y1_emu, y2_emu)))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(abs(x2_emu - x1_emu)))
    ext.set('cy', str(abs(y2_emu - y1_emu)))
    # Flip if needed
    if x2_emu < x1_emu:
        xfrm.set('flipH', '1')
    if y2_emu < y1_emu:
        xfrm.set('flipV', '1')

    prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
    prstGeom.set('prst', 'straightConnector1')
    etree.SubElement(prstGeom, qn('a:avLst'))

    return cxnSp


def _make_pic(x_emu, y_emu, w_emu, h_emu, rId, shape_id=None, name=None):
    """Create a <p:pic> element for an image."""
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    pic = etree.Element(qn('p:pic'), nsmap=nsmap)

    nvPicPr = etree.SubElement(pic, qn('p:nvPicPr'))
    cNvPr = etree.SubElement(nvPicPr, qn('p:cNvPr'))
    cNvPr.set('id', str(shape_id or 0))
    cNvPr.set('name', name or 'Picture')
    cNvPicPr = etree.SubElement(nvPicPr, qn('p:cNvPicPr'))
    etree.SubElement(cNvPicPr, qn('a:picLocks')).set('noChangeAspect', '1')
    etree.SubElement(nvPicPr, qn('p:nvPr'))

    blipFill = etree.SubElement(pic, qn('p:blipFill'))
    blip = etree.SubElement(blipFill, qn('a:blip'))
    blip.set(qn('r:embed'), rId)
    etree.SubElement(blipFill, qn('a:stretch')).append(etree.Element(qn('a:fillRect')))

    spPr = etree.SubElement(pic, qn('p:spPr'))
    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(x_emu))
    off.set('y', str(y_emu))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(w_emu))
    ext.set('cy', str(h_emu))
    prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
    prstGeom.set('prst', 'rect')
    etree.SubElement(prstGeom, qn('a:avLst'))

    return pic


def _make_pic_svg(x_emu, y_emu, w_emu, h_emu, rId, shape_id=None, name=None):
    """Create a <p:pic> element for an SVG image using asvg:svgBlip extension."""
    SVG_NS = 'http://schemas.microsoft.com/office/drawing/2016/SVG/main'
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'asvg': SVG_NS,
    }
    pic = etree.Element(qn('p:pic'), nsmap=nsmap)

    nvPicPr = etree.SubElement(pic, qn('p:nvPicPr'))
    cNvPr = etree.SubElement(nvPicPr, qn('p:cNvPr'))
    cNvPr.set('id', str(shape_id or 0))
    cNvPr.set('name', name or 'Picture')
    cNvPicPr = etree.SubElement(nvPicPr, qn('p:cNvPicPr'))
    etree.SubElement(cNvPicPr, qn('a:picLocks')).set('noChangeAspect', '1')
    etree.SubElement(nvPicPr, qn('p:nvPr'))

    blipFill = etree.SubElement(pic, qn('p:blipFill'))
    blipFill.set('rotWithShape', '1')
    blip = etree.SubElement(blipFill, qn('a:blip'))
    # SVG uses asvg:svgBlip extension instead of direct r:embed
    extLst = etree.SubElement(blip, qn('a:extLst'))
    ext = etree.SubElement(extLst, qn('a:ext'))
    ext.set('uri', '{96DAC541-7B7A-43D3-8B79-37D633B846F1}')
    svg_blip = etree.SubElement(ext, f'{{{SVG_NS}}}svgBlip')
    svg_blip.set(qn('r:embed'), rId)
    etree.SubElement(blipFill, qn('a:stretch'))

    spPr = etree.SubElement(pic, qn('p:spPr'))
    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(x_emu))
    off.set('y', str(y_emu))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(w_emu))
    ext.set('cy', str(h_emu))
    prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
    prstGeom.set('prst', 'rect')
    etree.SubElement(prstGeom, qn('a:avLst'))

    return pic


def _make_graphicFrame(x_emu, y_emu, w_emu, h_emu, tbl_element, shape_id=None, name=None):
    """Create a <p:graphicFrame> wrapping a table."""
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    gf = etree.Element(qn('p:graphicFrame'), nsmap=nsmap)

    nvGraphicFramePr = etree.SubElement(gf, qn('p:nvGraphicFramePr'))
    cNvPr = etree.SubElement(nvGraphicFramePr, qn('p:cNvPr'))
    cNvPr.set('id', str(shape_id or 0))
    cNvPr.set('name', name or 'Table')
    cNvGraphicFramePr = etree.SubElement(nvGraphicFramePr, qn('p:cNvGraphicFramePr'))
    etree.SubElement(cNvGraphicFramePr, qn('a:graphicFrameLocks')).set('noGrp', '1')
    etree.SubElement(nvGraphicFramePr, qn('p:nvPr'))

    xfrm = etree.SubElement(gf, qn('p:xfrm'))
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(x_emu))
    off.set('y', str(y_emu))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(w_emu))
    ext.set('cy', str(h_emu))

    graphic = etree.SubElement(gf, qn('a:graphic'))
    graphicData = etree.SubElement(graphic, qn('a:graphicData'))
    graphicData.set('uri', 'http://schemas.openxmlformats.org/drawingml/2006/table')
    graphicData.append(tbl_element)

    return gf


def _register_image(slide, image_path):
    """Add an image to the slide's package and return the relationship ID.

    Uses python-pptx's part management to add the image as a related part.
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    # Add image part to the slide part
    image_part = slide.part.package.get_or_add_image_part(image_path if hasattr(image_path, 'read') else str(image_path))
    rId = slide.part.relate_to(image_part, RT.IMAGE)
    return rId


def _register_svg(slide, svg_path):
    """Add an SVG file to the slide's package and return the relationship ID.

    SVG files can't be registered via get_or_add_image_part (PIL doesn't support SVG),
    so we create the part directly using python-pptx's Part class.
    """
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    svg_path = str(svg_path)
    with open(svg_path, 'rb') as f:
        svg_blob = f.read()

    # Find next available part name in the package
    pkg = slide.part.package
    max_num = 0
    for part in pkg.iter_parts():
        pn = str(part.partname)
        if pn.startswith('/ppt/media/image'):
            try:
                num = int(pn.split('image')[-1].split('.')[0])
                max_num = max(max_num, num)
            except ValueError:
                pass
    partname = PackURI(f'/ppt/media/image{max_num + 1}.svg')

    # Create the SVG part (Part(partname, content_type, package, blob))
    svg_part = Part(partname, 'image/svg+xml', pkg, svg_blob)

    # Create relationship from slide to SVG part
    rId = slide.part.relate_to(svg_part, RT.IMAGE)
    return rId


def _rgb(hex_color):
    """Convert hex color string to python-pptx RGBColor object."""
    tup = hex_to_rgb(hex_color)
    if tup is None:
        return None
    return RGBColor(*tup)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def resolve_string(value, theme):
    if not value:
        return value
    resolved, _ = resolve_theme_ref(value, theme)
    if isinstance(resolved, dict):
        return value
    return resolved


def build_base_style(content, theme):
    style = {}
    style_ref = content.get('style')
    if style_ref and isinstance(style_ref, str) and style_ref.startswith('$'):
        resolved, ok = resolve_theme_ref(style_ref, theme)
        if ok and isinstance(resolved, dict):
            style.update(resolved)
    elif style_ref and isinstance(style_ref, dict):
        style.update(style_ref)
    for key in ('fontSize', 'fontFamily', 'color', 'lineHeight', 'lineHeightPx',
                'letterSpacing', 'marginTop', 'textDirection', 'wrap', 'align'):
        if key in content:
            style[key] = content[key]
    if 'color' in style:
        style['color'] = resolve_string(style['color'], theme)
    return style


def _set_run_text_with_breaks(run, text):
    """Set run text, converting \\n to a:br elements for soft line breaks.

    In OOXML, <a:br> is a sibling of <a:r> at the paragraph level, not inside <a:r>.
    """
    if '\n' not in text:
        run.text = text
        return
    parts = text.split('\n')
    r_elem = run._r
    parent = r_elem.getparent()
    r_idx = list(parent).index(r_elem)
    # Copy the run's rPr for cloned runs
    rPr = r_elem.find(qn('a:rPr'))
    # Ensure source run has rPr (PowerPoint requires it)
    if rPr is None:
        rPr = etree.SubElement(r_elem, qn('a:rPr'))
        r_elem.insert(0, rPr)
    # Set first part as the run's text
    for t in r_elem.findall(qn('a:t')):
        r_elem.remove(t)
    t = etree.SubElement(r_elem, qn('a:t'))
    t.text = parts[0]
    t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
    # Insert a:br and new a:r for each subsequent part
    insert_idx = r_idx + 1
    for part in parts[1:]:
        br = etree.SubElement(parent, qn('a:br'))
        br.append(etree.fromstring(etree.tostring(rPr)))
        parent.insert(insert_idx, br)
        insert_idx += 1
        new_r = etree.SubElement(parent, qn('a:r'))
        new_r.append(etree.fromstring(etree.tostring(rPr)))
        new_t = etree.SubElement(new_r, qn('a:t'))
        new_t.text = part
        new_t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
        parent.insert(insert_idx, new_r)
        insert_idx += 1


def apply_text_style(run, style, theme):
    if 'font_size' in style:
        run.font.size = Pt(style['font_size'])
    elif 'fontSize' in style:
        run.font.size = Pt(style['fontSize'])

    if 'color' in style:
        color_val = style['color']
        rgb = resolve_color(color_val, theme)
        if rgb:
            run.font.color.rgb = _rgb(rgb)
        # Hex8 alpha
        if isinstance(color_val, str) and color_val.startswith('#'):
            hex_str = color_val.lstrip('#')
            if len(hex_str) >= 8:
                try:
                    alpha_8bit = int(hex_str[6:8], 16)
                    if alpha_8bit < 255:
                        rPr = run._r.get_or_add_rPr()
                        solid_fill = rPr.find(qn('a:solidFill'))
                        if solid_fill is None:
                            solid_fill = etree.SubElement(rPr, qn('a:solidFill'))
                        srgb_clr = solid_fill.find(qn('a:srgbClr'))
                        if srgb_clr is None:
                            srgb_clr = etree.SubElement(solid_fill, qn('a:srgbClr'))
                            srgb_clr.set('val', hex_str[0:6].upper())
                        alpha_elem = srgb_clr.find(qn('a:alpha'))
                        if alpha_elem is None:
                            alpha_elem = etree.SubElement(srgb_clr, qn('a:alpha'))
                        alpha_elem.set('val', str(int(alpha_8bit / 255 * 100000)))
                except ValueError:
                    pass

    font_str = style.get('font_family') or style.get('fontFamily')
    if font_str:
        fonts = [f.strip().strip('"') for f in font_str.split(',')]
        run.font.name = fonts[0]
        # CJK font (second entry)
        if len(fonts) > 1 and fonts[1]:
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', fonts[1])

    if style.get('bold'):
        run.font.bold = True
    if style.get('italic'):
        run.font.italic = True
    if style.get('underline'):
        run.font.underline = True
    if style.get('strikethrough'):
        run._r.get_or_add_rPr().set('strike', 'sngStrike')
    if style.get('superscript'):
        run.font.superscript = True
    if style.get('subscript'):
        run.font.subscript = True

    # Background color highlight (a:rPr/a:highlight/a:srgbClr)
    bg_color = style.get('bg_color')
    if bg_color:
        resolved_bg = resolve_color(bg_color, theme)
        if resolved_bg:
            rPr = run._r.get_or_add_rPr()
            for hl in rPr.findall(qn('a:highlight')):
                rPr.remove(hl)
            highlight = etree.SubElement(rPr, qn('a:highlight'))
            srgb = etree.SubElement(highlight, qn('a:srgbClr'))
            hex_str = resolved_bg.lstrip('#')
            if len(hex_str) >= 8:
                srgb.set('val', hex_str[:6])
                alpha_8bit = int(hex_str[6:8], 16)
                if alpha_8bit < 255:
                    alpha_pct = round(alpha_8bit / 255 * 100000)
                    alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                    alpha_elem.set('val', str(alpha_pct))
            else:
                srgb.set('val', hex_str)

    # Letter spacing (a:rPr/@spc, in 1/100 of a point)
    letter_spacing = style.get('letterSpacing')
    if letter_spacing:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(letter_spacing * 100)))

    # Hyperlink — store URL in rPr for post-processing
    hyperlink = style.get('hyperlink')
    if hyperlink:
        rPr = run._r.get_or_add_rPr()
        # Remove existing hlinkClick
        for hl in rPr.findall(qn('a:hlinkClick')):
            rPr.remove(hl)
        hlink = etree.SubElement(rPr, qn('a:hlinkClick'))
        hlink.set('action', 'ppaction://hlink')
        # Store URL as tooltip; will be converted to r:id by _resolve_hyperlinks
        hlink.set('tooltip', hyperlink)


def apply_alignment(paragraph, align_spec):
    if not align_spec:
        return
    if isinstance(align_spec, list) and len(align_spec) >= 1:
        h_align = align_spec[0]
    elif isinstance(align_spec, str):
        h_align = align_spec
    else:
        return
    m = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT,
         'justify': PP_ALIGN.JUSTIFY, 'distributed': PP_ALIGN.DISTRIBUTE}
    if h_align in m:
        paragraph.alignment = m[h_align]


def apply_vertical_alignment(text_frame, align_spec):
    if not align_spec:
        return
    if isinstance(align_spec, list) and len(align_spec) >= 2:
        v_align = align_spec[1]
    else:
        return
    m = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
    if v_align in m:
        text_frame.vertical_anchor = m[v_align]


def apply_paragraph_properties(paragraph, style):
    """Apply paragraph-level OOXML properties (spacing, indent) via direct XML."""
    pPr = paragraph._p.get_or_add_pPr()

    # Spacing before (marginTop in points → spcPts)
    margin_top = style.get('marginTop') or style.get('margin_top')
    if margin_top and margin_top > 0:
        spc_bef = etree.SubElement(pPr, qn('a:spcBef'))
        spc_pts = etree.SubElement(spc_bef, qn('a:spcPts'))
        spc_pts.set('val', str(int(margin_top * 100)))  # points * 100

    # Spacing after (marginBottom in points → spcPts)
    margin_bottom = style.get('marginBottom') or style.get('margin_bottom')
    if margin_bottom and margin_bottom > 0:
        spc_aft = etree.SubElement(pPr, qn('a:spcAft'))
        spc_pts = etree.SubElement(spc_aft, qn('a:spcPts'))
        spc_pts.set('val', str(int(margin_bottom * 100)))

    # Line height (lineHeight or lineHeightPx)
    line_height = style.get('lineHeight') or style.get('line_height')
    line_height_px = style.get('lineHeightPx')
    if line_height_px:
        # lineHeightPx is in pixels, convert to points (1px ≈ 0.75pt)
        line_height = line_height_px * 0.75
    if line_height:
        ln_spc = etree.SubElement(pPr, qn('a:lnSpc'))
        if isinstance(line_height, (int, float)) and line_height < 10:
            # It's a multiplier (e.g., 1.3)
            spc_pct = etree.SubElement(ln_spc, qn('a:spcPct'))
            spc_pct.set('val', str(int(line_height * 100000)))
        else:
            # It's in points
            spc_pts = etree.SubElement(ln_spc, qn('a:spcPts'))
            spc_pts.set('val', str(int(line_height * 100)))

    # Indentation (padding_left → marL, text_indent → firstLine)
    padding_left = style.get('padding_left')
    text_indent = style.get('text_indent')
    if padding_left or text_indent:
        indent = etree.SubElement(pPr, qn('a:indent'))
        if padding_left:
            # Convert points to EMU (1 pt = 12700 EMU)
            indent.set('marL', str(int(padding_left * 12700)))
        if text_indent:
            indent.set('firstLine', str(int(text_indent * 12700)))

    # Letter spacing — applied at run level via letterSpacing key
    # (letter_spacing from rich text <p> tags is converted to letterSpacing in build_base_style)


def _apply_list_formatting(paragraph, list_type, default_bullets, level=None):
    """Apply bullet/numbering formatting to a paragraph based on list_type.

    list_type is a list like ['ul'] or ['ul', 'ul'] indicating nesting depth.
    level is the OOXML level (0-based). If None, derived from len(list_type) - 1.
    """
    pPr = paragraph._p.get_or_add_pPr()
    if level is None:
        level = len(list_type) - 1  # 0-based nesting level
    pPr.set('lvl', str(level))

    last_type = list_type[-1]  # 'ul' or 'ol'
    if last_type == 'ul':
        # Unordered list — add buChar
        bullet_char = default_bullets[min(level, len(default_bullets) - 1)]
        bu_char = etree.SubElement(pPr, qn('a:buChar'))
        bu_char.set('char', bullet_char)
    elif last_type == 'ol':
        # Ordered list — add buAutoNum
        bu_auto = etree.SubElement(pPr, qn('a:buAutoNum'))
        bu_auto.set('type', 'arabicPeriod')


# ---------------------------------------------------------------------------
# OOXML fill injection (gradient, image fill)
# ---------------------------------------------------------------------------

def _inject_gradient_fill(shape, fill_spec, theme):
    """Inject a:gradFill directly into shape's spPr XML."""
    stops = fill_spec.get('stops', [])
    if not stops:
        return
    angle = fill_spec.get('angle', 0)
    gradient_type = fill_spec.get('gradientType', 'linear')

    spPr = shape._element.spPr
    # Remove existing fills
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:blipFill'):
        for child in spPr.findall(qn(tag)):
            spPr.remove(child)

    grad_fill = etree.SubElement(spPr, qn('a:gradFill'))
    gs_lst = etree.SubElement(grad_fill, qn('a:gsLst'))
    for stop in stops:
        pos = int(stop.get('position', 0) * 100000)
        color_hex = resolve_color(stop.get('color'), theme)
        if not color_hex:
            continue
        gs = etree.SubElement(gs_lst, qn('a:gs'))
        gs.set('pos', str(pos))
        color_clean = color_hex.lstrip('#')
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        # Check for alpha in hex string (8-char hex = #RRGGBBAA)
        if len(color_clean) > 6:
            srgb.set('val', color_clean[:6])
            alpha_8bit = int(color_clean[6:8], 16)
            alpha_ooxml = round(alpha_8bit / 255 * 100000)
            alpha_el = etree.SubElement(srgb, qn('a:alpha'))
            alpha_el.set('val', str(alpha_ooxml))
        else:
            srgb.set('val', color_clean)

    if gradient_type == 'radial':
        path = etree.SubElement(grad_fill, qn('a:path'))
        path.set('path', 'circle')
        fill_to_rect = etree.SubElement(path, qn('a:fillToRect'))
        fill_to_rect.set('l', '50000')
        fill_to_rect.set('t', '50000')
        fill_to_rect.set('r', '50000')
        fill_to_rect.set('b', '50000')
    else:
        lin = etree.SubElement(grad_fill, qn('a:lin'))
        lin.set('ang', str(int(angle * 60000)))
        lin.set('scaled', '1')

    # Insert after geometry (OOXML schema: xfrm, geom, fill, ln, ...)
    geom = spPr.find(qn('a:custGeom')) or spPr.find(qn('a:prstGeom'))
    if geom is not None:
        geom_idx = list(spPr).index(geom)
        spPr.remove(grad_fill)
        spPr.insert(geom_idx + 1, grad_fill)
    else:
        xfrm = spPr.find(qn('a:xfrm'))
        if xfrm is not None:
            xfrm_idx = list(spPr).index(xfrm)
            spPr.remove(grad_fill)
            spPr.insert(xfrm_idx + 1, grad_fill)
        else:
            spPr.remove(grad_fill)
            spPr.insert(0, grad_fill)


def _inject_image_fill(shape, fill_spec, theme, prs, base_dir):
    """Inject a:blipFill into shape's spPr."""
    src = fill_spec.get('src', '')
    if not src:
        return

    # Resolve image path
    img_path = None
    if src.startswith('http://') or src.startswith('https://'):
        try:
            import requests
            from io import BytesIO
            resp = requests.get(src, timeout=10)
            if resp.status_code == 200:
                img_path = BytesIO(resp.content)
        except Exception:
            return
    else:
        p = Path(src)
        if not p.is_absolute():
            p = base_dir / p
        if p.exists():
            img_path = str(p)

    if img_path is None:
        return

    # Register image and get relationship ID
    rId = _register_image(shape.part, img_path if isinstance(img_path, str) else img_path)

    spPr = shape._element.spPr
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:blipFill'):
        for child in spPr.findall(qn(tag)):
            spPr.remove(child)

    blip_fill = etree.SubElement(spPr, qn('a:blipFill'))
    blip = etree.SubElement(blip_fill, qn('a:blip'))
    blip.set(qn('r:embed'), rId)
    stretch = etree.SubElement(blip_fill, qn('a:stretch'))
    etree.SubElement(stretch, qn('a:fillRect'))

    # Insert after geometry (OOXML schema: xfrm, geom, fill, ln, ...)
    geom = spPr.find(qn('a:custGeom')) or spPr.find(qn('a:prstGeom'))
    if geom is not None:
        geom_idx = list(spPr).index(geom)
        spPr.remove(blip_fill)
        spPr.insert(geom_idx + 1, blip_fill)
    else:
        xfrm = spPr.find(qn('a:xfrm'))
        if xfrm is not None:
            xfrm_idx = list(spPr).index(xfrm)
            spPr.remove(blip_fill)
            spPr.insert(xfrm_idx + 1, blip_fill)
        else:
            spPr.remove(blip_fill)
            spPr.insert(0, blip_fill)


def _inject_shadow(shape, shadow_spec, theme):
    """Inject a:effectLst/a:outerShdw into shape XML."""
    blur = shadow_spec.get('blur', 0)
    offset = shadow_spec.get('offset', [0, 0])
    color = resolve_color(shadow_spec.get('color'), theme) or '#000000'

    sp = shape._element
    spPr = sp.spPr

    # Remove existing effectLst
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)

    effect_lst = etree.SubElement(spPr, qn('a:effectLst'))
    outer = etree.SubElement(effect_lst, qn('a:outerShdw'))
    outer.set('blurRad', str(px_to_emu(blur)))
    dist = math.sqrt(offset[0]**2 + offset[1]**2) if offset else 0
    outer.set('dist', str(px_to_emu(int(dist))))
    if dist > 0:
        angle = math.atan2(offset[1], offset[0])
        outer.set('dir', str(int(angle * 60000 * 180 / math.pi)))
    outer.set('algn', 'tl')
    outer.set('rotWithShape', '0')
    srgb = etree.SubElement(outer, qn('a:srgbClr'))
    hex_str = color.lstrip('#')
    if len(hex_str) >= 8:
        srgb.set('val', hex_str[:6])
        alpha_8bit = int(hex_str[6:8], 16)
        if alpha_8bit < 255:
            alpha_pct = round(alpha_8bit / 255 * 100000)
            alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
            alpha_elem.set('val', str(alpha_pct))
    else:
        srgb.set('val', hex_str)


def apply_fill(shape, fill_spec, theme, prs=None, base_dir=None):
    """Apply fill to a shape, handling solid/gradient/image/noFill."""
    if fill_spec is None:
        # Explicit noFill — inject <a:noFill/> into spPr
        try:
            spPr = shape._element.spPr
            # Remove existing fill elements
            for tag in ('a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:noFill'):
                for el in spPr.findall(qn(tag)):
                    spPr.remove(el)
            # Insert noFill at beginning (before a:ln if present)
            no_fill = etree.SubElement(spPr, qn('a:noFill'))
            ln = spPr.find(qn('a:ln'))
            if ln is not None:
                spPr.remove(no_fill)
                spPr.insert(list(spPr).index(ln), no_fill)
        except Exception:
            pass
        return
    if not fill_spec:
        return
    fill_type = fill_spec.get('type')
    if fill_type == 'solid':
        color_val = fill_spec.get('color')
        rgb = resolve_color(color_val, theme)
        if rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(rgb)
            # Handle alpha in hex8 colors via direct XML injection
            color_clean = rgb.lstrip('#')
            if len(color_clean) >= 8:
                alpha_8bit = int(color_clean[6:8], 16)
                if alpha_8bit < 255:
                    spPr = shape._element.spPr
                    sf = spPr.find(qn('a:solidFill'))
                    if sf is not None:
                        srgb = sf.find(qn('a:srgbClr'))
                        if srgb is not None:
                            srgb.set('val', color_clean[:6].upper())
                            alpha_elem = srgb.find(qn('a:alpha'))
                            if alpha_elem is None:
                                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                            alpha_elem.set('val', str(int(alpha_8bit / 255 * 100000)))
    elif fill_type == 'gradient':
        _inject_gradient_fill(shape, fill_spec, theme)
    elif fill_type == 'image':
        if prs:
            _inject_image_fill(shape, fill_spec, theme, prs, base_dir)


def apply_border(shape, border_spec, theme):
    if not border_spec:
        return
    style = border_spec.get('style', 'solid')
    if style == 'none':
        shape.line.fill.background()
        return
    width = border_spec.get('width', 1)
    color_val = border_spec.get('color')
    rgb = resolve_color(color_val, theme)
    if rgb:
        shape.line.color.rgb = _rgb(rgb)
        # Handle alpha in hex8 colors via direct XML injection
        color_clean = rgb.lstrip('#')
        if len(color_clean) > 6:
            ln = shape._element.spPr.find(qn('a:ln'))
            if ln is not None:
                sf = ln.find(qn('a:solidFill'))
                if sf is not None:
                    srgb = sf.find(qn('a:srgbClr'))
                    if srgb is not None:
                        srgb.set('val', color_clean[:6])
                        alpha_8bit = int(color_clean[6:8], 16)
                        alpha_ooxml = round(alpha_8bit / 255 * 100000)
                        alpha_el = etree.SubElement(srgb, qn('a:alpha'))
                        alpha_el.set('val', str(alpha_ooxml))
    shape.line.width = Pt(width)
    if style == 'dash':
        shape.line.dash_style = 2
    elif style == 'dot':
        shape.line.dash_style = 3
    else:
        shape.line.dash_style = 1


def apply_transform(shape, element):
    """Apply rotation and flip from element to shape."""
    rotation = element.get('rotation')
    if rotation:
        shape.rotation = rotation
    flip = element.get('flip')
    if flip and isinstance(flip, list) and len(flip) >= 2:
        xfrm = shape._element.spPr.find(qn('a:xfrm'))
        if xfrm is not None:
            if flip[0]:
                xfrm.set('flipH', '1')
            if flip[1]:
                xfrm.set('flipV', '1')


def apply_opacity(shape, opacity):
    """Apply element-level opacity via a:effectLst/a:alpha (OOXML standard)."""
    if opacity is None or opacity >= 1:
        return
    spPr = shape._element.spPr
    effect_lst = spPr.find(qn('a:effectLst'))
    if effect_lst is None:
        effect_lst = etree.SubElement(spPr, qn('a:effectLst'))
    alpha_elem = effect_lst.find(qn('a:alpha'))
    if alpha_elem is None:
        alpha_elem = etree.SubElement(effect_lst, qn('a:alpha'))
    alpha_elem.set('val', str(int(opacity * 100000)))


# ---------------------------------------------------------------------------
# Element renderers
# ---------------------------------------------------------------------------

def add_text_element(slide, element, theme, prs=None, base_dir=None):
    bounds = element.get('bounds', [0, 0, 100, 100])
    x, y, w, h = bounds

    # Try to use layout placeholder if element has placeholder metadata
    ph_meta = element.get('placeholder')
    shape = None
    is_placeholder = False
    layout_paras = None  # Layout's paragraph pPr elements for bullet inheritance
    if ph_meta:
        ph_idx = ph_meta.get('idx')
        if ph_idx is not None:
            # First try slide.placeholders (only title/body are inherited by python-pptx)
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == ph_idx:
                    shape = ph
                    is_placeholder = True
                    # Update position/size from element data
                    shape.left = px_to_emu(x)
                    shape.top = px_to_emu(y)
                    shape.width = px_to_emu(w)
                    shape.height = px_to_emu(h)
                    # Get layout's paragraph pPr for bullet/indent inheritance
                    layout_ph = _find_layout_placeholder(slide, ph_idx)
                    if layout_ph is not None and layout_ph.has_text_frame:
                        layout_paras = [p._p for p in layout_ph.text_frame.paragraphs]
                    break

            # If not found on slide, try copying from layout (footer/sldNum/date/etc.)
            if shape is None:
                layout_ph = _find_layout_placeholder(slide, ph_idx)
                if layout_ph is not None:
                    from copy import deepcopy
                    sp_copy = deepcopy(layout_ph._element)
                    sp_tree = slide._element.find(qn('p:cSld') + '/' + qn('p:spTree'))
                    if sp_tree is None:
                        sp_tree = slide._element.find(qn('p:spTree'))
                    if sp_tree is not None:
                        sp_tree.append(sp_copy)
                        # Find the newly added shape
                        for ph in slide.placeholders:
                            if ph.placeholder_format.idx == ph_idx:
                                shape = ph
                                is_placeholder = True
                                shape.left = px_to_emu(x)
                                shape.top = px_to_emu(y)
                                shape.width = px_to_emu(w)
                                shape.height = px_to_emu(h)
                                if layout_ph.has_text_frame:
                                    layout_paras = [p._p for p in layout_ph.text_frame.paragraphs]
                                break

    if shape is None:
        # Build shape directly as XML — no python-pptx defaults
        shape_name = element.get('shapeName', 'rect')
        # Use shapeName for display name
        display_name = shape_name[0].upper() + shape_name[1:] if shape_name else 'Shape'
        sp = _make_sp(px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h),
                      shape_name=shape_name, txBox=True,
                      name=display_name,
                      shape_id=_next_shape_id(slide))
        sp_tree = _get_sp_tree(slide)
        sp_tree.append(sp)
        # Find the shape via slide.shapes to get the python-pptx wrapper
        for s in slide.shapes:
            if s._element is sp:
                shape = s
                break
        if shape is None:
            shape = sp  # fallback to raw element
    tf = shape.text_frame
    if not is_placeholder:
        # Only set word_wrap for non-placeholder shapes.
        # Placeholders inherit bodyPr from layout/master; modifying it breaks inheritance.
        tf.word_wrap = element.get('content', {}).get('wrap', True)

    # Text direction (vertical text)
    text_direction = element.get('content', {}).get('textDirection')
    if text_direction == 'vertical':
        body_pr = tf._txBody.find(qn('a:bodyPr'))
        if body_pr is not None:
            body_pr.set('vert', 'vert')

    content = element.get('content', {})
    base_style = build_base_style(content, theme)
    text = content.get('text', '')
    segments = parse_rich_text(text, base_style)

    if not segments:
        return shape

    align_spec = content.get('align')
    if align_spec:
        apply_vertical_alignment(tf, align_spec)

    if is_placeholder:
        _update_placeholder_text(tf, segments, base_style, theme, align_spec, layout_paras)
    else:
        _create_new_text(tf, segments, base_style, theme, align_spec)

    # Text gradient
    gradient = content.get('gradient')
    if gradient:
        for para in tf.paragraphs:
            for run in para.runs:
                _inject_text_gradient(run, gradient, theme)

    # Text shadow
    shadow = content.get('shadow')
    if shadow:
        for para in tf.paragraphs:
            for run in para.runs:
                _inject_text_shadow(run, shadow, theme)

    # Apply fill and border (for shapes with background colors)
    fill_spec = element.get('fill')
    if fill_spec:
        apply_fill(shape, fill_spec, theme, prs, base_dir)
    border_spec = element.get('border')
    if border_spec:
        apply_border(shape, border_spec, theme)

    # Apply element-level transform
    apply_transform(shape, element)
    opacity = element.get('opacity')
    if opacity is not None and opacity < 1:
        apply_opacity(shape, opacity)

    # Resolve hyperlinks (convert tooltip URLs to r:id relationships)
    try:
        _resolve_hyperlinks(slide, shape)
    except Exception:
        pass  # Hyperlink resolution is best-effort

    return shape


def _resolve_hyperlinks(slide, shape):
    """Convert tooltip-stored URLs in hlinkClick to proper r:id relationships."""
    try:
        tf = shape.text_frame
    except AttributeError:
        return

    # Collect all hyperlinks first, then create relationships
    hyperlinks = []  # list of (hlink_element, url)
    for para in tf.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is None:
                continue
            for hlink in rPr.findall(qn('a:hlinkClick')):
                url = hlink.get('tooltip')
                if url:
                    hyperlinks.append((hlink, url))

    if not hyperlinks:
        return

    # Create relationships for all hyperlinks
    # Get the slide's relationships part
    rels_part = slide.part._rels
    for hlink, url in hyperlinks:
        # Generate a unique rId
        rId = f'rId{len(rels_part) + 100}'  # Offset to avoid conflicts
        while rId in rels_part:
            rId = f'rId{int(rId[4:]) + 1}'

        # Create the relationship element directly
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        rels_part._rels[rId] = _ExternalHyperlinkRel(rId, url, RT.HYPERLINK)
        hlink.set(qn('r:id'), rId)
        del hlink.attrib['tooltip']


class _ExternalHyperlinkRel:
    """Minimal relationship object for external hyperlinks."""
    def __init__(self, rId, target, rel_type):
        self.rId = rId
        self._target = target
        self.reltype = rel_type
        self.is_external = True
    @property
    def target_ref(self):
        return self._target


def _find_layout_placeholder(slide, ph_idx):
    """Find the placeholder in the slide's layout (not the slide itself)."""
    try:
        layout = slide.slide_layout
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == ph_idx:
                return ph
    except Exception:
        pass
    return None


def _group_segments_by_para(segments):
    """Group segments into paragraphs. Returns list of (style, [segments])."""
    paras = []
    current_style = {}
    current_segs = []
    for seg in segments:
        if seg.get('is_new_para', False) or not current_segs:
            if current_segs:
                paras.append((current_style, current_segs))
            current_style = seg.get('style', {})
            current_segs = [seg]
        else:
            current_segs.append(seg)
    if current_segs:
        paras.append((current_style, current_segs))
    return paras


def _update_placeholder_text(tf, segments, base_style, theme, align_spec,
                             layout_paras=None):
    """Update placeholder text while preserving layout-inherited formatting.

    Instead of replacing the entire text frame (which destroys bullets, spacing,
    etc. inherited from the layout), we update text within existing paragraphs.
    For existing paragraphs, we clear only the a:t text nodes and update them,
    preserving a:rPr formatting and a:pPr (which controls bullets via lvl attr).
    This maintains OOXML inheritance from the layout/master body style.
    """
    para_groups = _group_segments_by_para(segments)
    existing_paras = list(tf.paragraphs)
    _DEFAULT_BULLETS = ['\u2022']  # Bullet character (•) — level only affects indent, not symbol

    for i, (para_style, segs) in enumerate(para_groups):
        if i < len(existing_paras):
            # Update existing paragraph in-place to preserve pPr and rPr
            para = existing_paras[i]
            p_elem = para._p

            # Copy pPr from layout if paragraph doesn't have one
            if layout_paras and i < len(layout_paras):
                existing_pPr = p_elem.find(qn('a:pPr'))
                if existing_pPr is None:
                    _copy_pPr(layout_paras[i], p_elem)

            # Apply bullet formatting for list items, or clear layout-inherited bullets
            first_seg = segs[0] if segs else {}
            list_type = first_seg.get('list_type')
            if list_type:
                pPr = p_elem.find(qn('a:pPr'))
                if pPr is None:
                    pPr = etree.SubElement(p_elem, qn('a:pPr'))
                    p_elem.insert(0, pPr)
                # Use list_level from PPTD if available, otherwise derive from nesting
                seg_level = first_seg.get('list_level')
                level = int(seg_level) if seg_level is not None and seg_level != '' else len(list_type) - 1
                pPr.set('lvl', str(level))
                if pPr.find(qn('a:buChar')) is None and pPr.find(qn('a:buNone')) is None:
                    last_type = list_type[-1]
                    if last_type == 'ul':
                        bullet_char = _DEFAULT_BULLETS[min(level, len(_DEFAULT_BULLETS) - 1)]
                        bu_char = etree.SubElement(pPr, qn('a:buChar'))
                        bu_char.set('char', bullet_char)
                    elif last_type == 'ol':
                        bu_auto = etree.SubElement(pPr, qn('a:buAutoNum'))
                        bu_auto.set('type', 'arabicPeriod')
            else:
                # Not a list item — remove layout-inherited lvl and bullet elements
                pPr = p_elem.find(qn('a:pPr'))
                if pPr is not None:
                    if 'lvl' in pPr.attrib:
                        del pPr.attrib['lvl']
                    for bu_tag in ('a:buChar', 'a:buAutoNum', 'a:buNone', 'a:buFont', 'a:buSzPts', 'a:buClr'):
                        for bu in pPr.findall(qn(bu_tag)):
                            pPr.remove(bu)

            # Collect existing run elements
            run_elems = p_elem.findall(qn('a:r'))

            # Clear text from existing runs (preserve rPr)
            for r_elem in run_elems:
                t_elem = r_elem.find(qn('a:t'))
                if t_elem is not None:
                    t_elem.text = ''

            # Also handle endParaRPr — convert it to a run if we need more runs
            end_para = p_elem.find(qn('a:endParaRPr'))

            # Now assign text to runs
            seg_idx = 0
            for r_idx, r_elem in enumerate(run_elems):
                if seg_idx < len(segs):
                    t_elem = r_elem.find(qn('a:t'))
                    if t_elem is not None:
                        seg_text = segs[seg_idx]['text']
                        if '\n' in seg_text:
                            # Remove the a:t and replace with a:t + a:br sequence
                            r_elem.remove(t_elem)
                            parts = seg_text.split('\n')
                            new_t = etree.SubElement(r_elem, qn('a:t'))
                            new_t.text = parts[0]
                            new_t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
                            # Insert a:br and new a:r as siblings after this run
                            rPr_copy = r_elem.find(qn('a:rPr'))
                            if rPr_copy is None:
                                rPr_copy = etree.SubElement(r_elem, qn('a:rPr'))
                            parent = r_elem.getparent()
                            insert_idx = list(parent).index(r_elem) + 1
                            for part in parts[1:]:
                                br = etree.SubElement(parent, qn('a:br'))
                                if rPr_copy is not None:
                                    br.append(etree.fromstring(etree.tostring(rPr_copy)))
                                parent.insert(insert_idx, br)
                                insert_idx += 1
                                new_r = etree.SubElement(parent, qn('a:r'))
                                if rPr_copy is not None:
                                    new_r.append(etree.fromstring(etree.tostring(rPr_copy)))
                                br_t = etree.SubElement(new_r, qn('a:t'))
                                br_t.text = part
                                br_t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
                                parent.insert(insert_idx, new_r)
                                insert_idx += 1
                        else:
                            t_elem.text = seg_text
                        # Apply run-level style (font, color, etc.)
                        seg_style = segs[seg_idx].get('style', {})
                        merged = dict(base_style)
                        merged.update(seg_style)
                        # Update rPr if we have style info
                        rPr = r_elem.find(qn('a:rPr'))
                        if rPr is None:
                            rPr = etree.SubElement(r_elem, qn('a:rPr'))
                            r_elem.insert(0, rPr)
                        _apply_rpr_style(rPr, merged, theme)
                    seg_idx += 1

            # If we have more segments than existing runs, add new runs
            while seg_idx < len(segs):
                # Remove endParaRPr if present (it must come last)
                if end_para is not None:
                    p_elem.remove(end_para)
                    end_para = None
                run = para.add_run()
                _set_run_text_with_breaks(run, segs[seg_idx]['text'])
                seg_style = segs[seg_idx].get('style', {})
                merged = dict(base_style)
                merged.update(seg_style)
                apply_text_style(run, merged, theme)
                run._r.get_or_add_rPr()
                seg_idx += 1

        else:
            # Add new paragraph (beyond existing ones)
            para = tf.add_paragraph()
            # Copy pPr from layout paragraph if available (for bullet inheritance)
            if layout_paras and i < len(layout_paras):
                _copy_pPr(layout_paras[i], para._p)
            apply_paragraph_properties(para, para_style)

            # Apply bullet formatting for list items, or clear layout-inherited bullets
            first_seg = segs[0] if segs else {}
            list_type = first_seg.get('list_type')
            p_elem = para._p
            if list_type:
                pPr = p_elem.find(qn('a:pPr'))
                if pPr is None:
                    pPr = etree.SubElement(p_elem, qn('a:pPr'))
                    p_elem.insert(0, pPr)
                seg_level = first_seg.get('list_level')
                level = int(seg_level) if seg_level is not None and seg_level != '' else len(list_type) - 1
                pPr.set('lvl', str(level))
                if pPr.find(qn('a:buChar')) is None and pPr.find(qn('a:buNone')) is None:
                    last_type = list_type[-1]
                    if last_type == 'ul':
                        bullet_char = _DEFAULT_BULLETS[min(level, len(_DEFAULT_BULLETS) - 1)]
                        bu_char = etree.SubElement(pPr, qn('a:buChar'))
                        bu_char.set('char', bullet_char)
                    elif last_type == 'ol':
                        bu_auto = etree.SubElement(pPr, qn('a:buAutoNum'))
                        bu_auto.set('type', 'arabicPeriod')
            else:
                # Not a list item — remove layout-inherited lvl and bullet elements
                pPr = p_elem.find(qn('a:pPr'))
                if pPr is not None:
                    if 'lvl' in pPr.attrib:
                        del pPr.attrib['lvl']
                    for bu_tag in ('a:buChar', 'a:buAutoNum', 'a:buNone', 'a:buFont', 'a:buSzPts', 'a:buClr'):
                        for bu in pPr.findall(qn(bu_tag)):
                            pPr.remove(bu)
            for seg in segs:
                run = para.add_run()
                _set_run_text_with_breaks(run, seg['text'])
                seg_style = seg.get('style', {})
                merged = dict(base_style)
                merged.update(seg_style)
                apply_text_style(run, merged, theme)
                run._r.get_or_add_rPr()

        # Apply alignment
        if 'align' in para_style:
            apply_alignment(para, para_style['align'])
        elif align_spec:
            apply_alignment(para, align_spec)


def _copy_pPr(src_p, dst_p):
    """Copy a:pPr element from source paragraph to destination paragraph.

    This preserves bullet formatting (lvl attribute) and other paragraph properties
    that are inherited from the slide layout/master body style.
    """
    src_pPr = src_p.find(qn('a:pPr'))
    if src_pPr is None:
        return
    dst_pPr = dst_p.find(qn('a:pPr'))
    if dst_pPr is not None:
        dst_p.remove(dst_pPr)
    # Deep copy the pPr element
    from copy import deepcopy
    new_pPr = deepcopy(src_pPr)
    # Insert pPr as the first child of the paragraph
    dst_p.insert(0, new_pPr)


def _apply_rpr_style(rPr, style, theme):
    """Apply run-level formatting to an existing a:rPr element."""
    # Font size (camelCase to match PPTD format)
    font_size = style.get('fontSize')
    if font_size:
        rPr.set('sz', str(int(font_size * 100)))

    # Bold
    if style.get('bold'):
        rPr.set('b', '1')

    # Italic
    if style.get('italic'):
        rPr.set('i', '1')

    # Underline
    if style.get('underline'):
        rPr.set('u', 'sng')

    # Strikethrough
    if style.get('strikethrough'):
        rPr.set('strike', 'sngStrike')

    # Superscript / Subscript
    if style.get('superscript'):
        rPr.set('baseline', '30000')
    elif style.get('subscript'):
        rPr.set('baseline', '-25000')

    # Color
    color = style.get('color')
    if color:
        resolved = resolve_color(color, theme)
        if resolved:
            for sf in rPr.findall(qn('a:solidFill')):
                rPr.remove(sf)
            solid = etree.SubElement(rPr, qn('a:solidFill'))
            srgb = etree.SubElement(solid, qn('a:srgbClr'))
            hex_str = resolved.lstrip('#')
            if len(hex_str) >= 8:
                # Hex8: RRGGBBAA — extract alpha
                srgb.set('val', hex_str[:6])
                alpha_8bit = int(hex_str[6:8], 16)
                if alpha_8bit < 255:
                    alpha_pct = round(alpha_8bit / 255 * 100000)
                    alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                    alpha_elem.set('val', str(alpha_pct))
            else:
                srgb.set('val', hex_str)

    # Font family
    font_family = style.get('fontFamily') or style.get('font_family')
    if font_family:
        fonts = [f.strip().strip('"') for f in font_family.split(',')]
        latins = rPr.findall(qn('a:latin'))
        if latins:
            for latin in latins:
                latin.set('typeface', fonts[0])
        else:
            latin = etree.SubElement(rPr, qn('a:latin'))
            latin.set('typeface', fonts[0])
        # CJK font (second entry)
        if len(fonts) > 1 and fonts[1]:
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', fonts[1])

    # Letter spacing (a:rPr/@spc, in 1/100 of a point)
    letter_spacing = style.get('letterSpacing')
    if letter_spacing:
        rPr.set('spc', str(int(letter_spacing * 100)))

    # Background color highlight (a:rPr/a:highlight/a:srgbClr)
    bg_color = style.get('bg_color')
    if bg_color:
        resolved_bg = resolve_color(bg_color, theme)
        if resolved_bg:
            for hl in rPr.findall(qn('a:highlight')):
                rPr.remove(hl)
            highlight = etree.SubElement(rPr, qn('a:highlight'))
            srgb = etree.SubElement(highlight, qn('a:srgbClr'))
            hex_str = resolved_bg.lstrip('#')
            if len(hex_str) >= 8:
                srgb.set('val', hex_str[:6])
                alpha_8bit = int(hex_str[6:8], 16)
                if alpha_8bit < 255:
                    alpha_pct = round(alpha_8bit / 255 * 100000)
                    alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                    alpha_elem.set('val', str(alpha_pct))
            else:
                srgb.set('val', hex_str)

    # Hyperlink — store URL in rPr for post-processing
    hyperlink = style.get('hyperlink')
    if hyperlink:
        for hl in rPr.findall(qn('a:hlinkClick')):
            rPr.remove(hl)
        hlink = etree.SubElement(rPr, qn('a:hlinkClick'))
        hlink.set('action', 'ppaction://hlink')
        hlink.set('tooltip', hyperlink)


def _create_new_text(tf, segments, base_style, theme, align_spec):
    """Create new text paragraphs (for non-placeholder text boxes)."""
    current_para = None
    para_idx = 0

    # Default bullet characters by nesting level
    _DEFAULT_BULLETS = ['\u2022']  # Bullet character (•) — level only affects indent, not symbol

    for seg in segments:
        is_new_para = seg.get('is_new_para', False) or para_idx == 0
        if is_new_para or current_para is None:
            current_para = tf.paragraphs[0] if para_idx == 0 else tf.add_paragraph()
            para_idx += 1

            # Apply paragraph-level properties (spacing, indent)
            seg_style_for_para = seg.get('style', {})
            apply_paragraph_properties(current_para, seg_style_for_para)

            # Apply bullet formatting if segment is a list item
            list_type = seg.get('list_type')
            if list_type:
                seg_level = seg.get('list_level')
                level = int(seg_level) if seg_level is not None and seg_level != '' else None
                _apply_list_formatting(current_para, list_type, _DEFAULT_BULLETS, level=level)

        run = current_para.add_run()
        _set_run_text_with_breaks(run, seg['text'])

        seg_style = seg.get('style', {})
        merged_style = dict(base_style)
        merged_style.update(seg_style)
        apply_text_style(run, merged_style, theme)
        # Ensure every run has an rPr (PowerPoint requires it)
        run._r.get_or_add_rPr()

        if 'align' in seg_style:
            apply_alignment(current_para, seg_style['align'])
        elif align_spec:
            apply_alignment(current_para, align_spec)


def _inject_text_gradient(run, gradient_spec, theme):
    """Inject gradient fill into text run's rPr."""
    stops = gradient_spec.get('stops', [])
    if not stops:
        return
    rPr = run._r.get_or_add_rPr()
    for sf in rPr.findall(qn('a:solidFill')):
        rPr.remove(sf)
    grad_fill = etree.SubElement(rPr, qn('a:gradFill'))
    gs_lst = etree.SubElement(grad_fill, qn('a:gsLst'))
    for stop in stops:
        pos = int(stop.get('position', 0) * 100000)
        color_hex = resolve_color(stop.get('color'), theme)
        if not color_hex:
            continue
        gs = etree.SubElement(gs_lst, qn('a:gs'))
        gs.set('pos', str(pos))
        color_clean = color_hex.lstrip('#')
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        # Check for alpha in hex string (8-char hex = #RRGGBBAA)
        if len(color_clean) > 6:
            srgb.set('val', color_clean[:6])
            alpha_8bit = int(color_clean[6:8], 16)
            alpha_ooxml = round(alpha_8bit / 255 * 100000)
            alpha_el = etree.SubElement(srgb, qn('a:alpha'))
            alpha_el.set('val', str(alpha_ooxml))
        else:
            srgb.set('val', color_clean)
    angle = gradient_spec.get('angle', 0)
    lin = etree.SubElement(grad_fill, qn('a:lin'))
    lin.set('ang', str(int(angle * 60000)))
    lin.set('scaled', '1')


def _inject_text_shadow(run, shadow_spec, theme):
    """Inject shadow into text run's rPr."""
    blur = shadow_spec.get('blur', 0)
    offset = shadow_spec.get('offset', [0, 0])
    color = resolve_color(shadow_spec.get('color'), theme) or '#000000'
    rPr = run._r.get_or_add_rPr()
    for el in rPr.findall(qn('a:effectLst')):
        rPr.remove(el)
    effect_lst = etree.SubElement(rPr, qn('a:effectLst'))
    outer = etree.SubElement(effect_lst, qn('a:outerShdw'))
    outer.set('blurRad', str(px_to_emu(blur)))
    dist = math.sqrt(offset[0]**2 + offset[1]**2) if offset else 0
    outer.set('dist', str(px_to_emu(int(dist))))
    if dist > 0:
        angle = math.atan2(offset[1], offset[0])
        outer.set('dir', str(int(angle * 60000 * 180 / math.pi)))
    outer.set('algn', 'tl')
    outer.set('rotWithShape', '0')
    srgb = etree.SubElement(outer, qn('a:srgbClr'))
    hex_str = color.lstrip('#')
    if len(hex_str) >= 8:
        srgb.set('val', hex_str[:6])
        alpha_8bit = int(hex_str[6:8], 16)
        if alpha_8bit < 255:
            alpha_pct = round(alpha_8bit / 255 * 100000)
            alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
            alpha_elem.set('val', str(alpha_pct))
    else:
        srgb.set('val', hex_str)


def _apply_arrows(shape, arrow):
    """Apply arrow endpoints to a line/connector shape via a:ln/a:headEnd and a:tailEnd."""
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))

    # arrow[0] = start arrow, arrow[1] = end arrow
    if len(arrow) >= 1 and arrow[0]:
        head_end = ln.find(qn('a:headEnd'))
        if head_end is None:
            head_end = etree.SubElement(ln, qn('a:headEnd'))
        head_end.set('type', arrow[0])

    if len(arrow) >= 2 and arrow[1]:
        tail_end = ln.find(qn('a:tailEnd'))
        if tail_end is None:
            tail_end = etree.SubElement(ln, qn('a:tailEnd'))
        tail_end.set('type', arrow[1])


def add_shape_element(slide, element, theme, prs=None, base_dir=None):
    bounds = element.get('bounds', [0, 0, 100, 100])
    x, y, w, h = bounds
    shape_name = element.get('shapeName', 'rect')

    # Line shapes
    line_shapes = {'straightConnector1', 'bentConnector2', 'bentConnector3',
                   'bentConnector4', 'curvedConnector2', 'curvedConnector3',
                   'curvedConnector4', 'line'}
    is_line = shape_name in line_shapes

    if is_line:
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        # Use actual connector shape instead of rectangle
        begin_x, begin_y = px_to_emu(x), px_to_emu(y)
        end_x = px_to_emu(x + max(w, 1))
        end_y = px_to_emu(y + max(h, 1))
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, begin_x, begin_y, end_x, end_y)
        border_spec = element.get('border')
        if border_spec:
            apply_border(shape, border_spec, theme)
        # Arrows
        arrow = element.get('arrow')
        if arrow and isinstance(arrow, list):
            _apply_arrows(shape, arrow)
    elif shape_name == 'custom':
        # Custom path — try to build custGeom, fallback to rect
        path_data = element.get('path')
        if path_data:
            shape = _add_custom_path(slide, element, path_data, x, y, w, h, theme)
        else:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h))
            # Remove <p:style> — it adds unwanted theme fill/line references
            style_el = shape._element.find(qn('p:style'))
            if style_el is not None:
                shape._element.remove(style_el)
        fill_spec = element.get('fill')
        if fill_spec:
            apply_fill(shape, fill_spec, theme, prs, base_dir)
        border_spec = element.get('border')
        if border_spec:
            apply_border(shape, border_spec, theme)
    else:
        mso_map = load_pptd_to_mso_map()
        mso = mso_map.get(shape_name, MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(mso, px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h))

        # Remove <p:style> — it adds unwanted theme fill/line references
        style_el = shape._element.find(qn('p:style'))
        if style_el is not None:
            shape._element.remove(style_el)

        adjustments = element.get('adjustments')
        if adjustments and hasattr(shape, 'adjustments'):
            try:
                for i, adj in enumerate(adjustments):
                    if i < len(shape.adjustments):
                        shape.adjustments[i] = adj / 100000.0
            except Exception:
                pass

        fill_spec = element.get('fill')
        if fill_spec:
            apply_fill(shape, fill_spec, theme, prs, base_dir)
        border_spec = element.get('border')
        if border_spec:
            apply_border(shape, border_spec, theme)

    # Shadow
    shadow = element.get('shadow')
    if shadow:
        _inject_shadow(shape, shadow, theme)

    apply_transform(shape, element)
    opacity = element.get('opacity')
    if opacity is not None and opacity < 1:
        apply_opacity(shape, opacity)

    return shape


def _add_custom_path(slide, element, path_data, x, y, w, h, theme):
    """Create a shape with custom geometry from PPTD path string."""
    # Path format: "viewBoxW,viewBoxH;SVG_path_data"
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h))
    # Remove <p:style> — it adds unwanted theme fill/line references
    style_el = shape._element.find(qn('p:style'))
    if style_el is not None:
        shape._element.remove(style_el)
    try:
        parts = path_data.split(';', 1)
        if len(parts) != 2:
            return shape
        viewbox = parts[0].split(',')
        svg_path = parts[1]
        vw, vh = int(viewbox[0]), int(viewbox[1])

        # Set shape size from page bounds (not viewBox) — viewBox defines the
        # path's coordinate system, shape ext defines the display size.
        # PowerPoint scales the path to fit within the shape bounds.
        xfrm = shape._element.spPr.find(qn('a:xfrm'))
        if xfrm is not None:
            ext = xfrm.find(qn('a:ext'))
            if ext is not None:
                ext.set('cx', str(px_to_emu(w)))
                ext.set('cy', str(px_to_emu(h)))

        spPr = shape._element.spPr
        # Remove preset geometry
        for pg in spPr.findall(qn('a:prstGeom')):
            spPr.remove(pg)
        # Remove rect geometry
        for rg in spPr.findall(qn('a:custGeom')):
            spPr.remove(rg)

        cust_geom = etree.SubElement(spPr, qn('a:custGeom'))
        etree.SubElement(cust_geom, qn('a:avLst'))
        etree.SubElement(cust_geom, qn('a:gdLst'))
        etree.SubElement(cust_geom, qn('a:ahLst'))
        etree.SubElement(cust_geom, qn('a:cxnLst'))
        # Text rectangle defaults to full shape area
        rect = etree.SubElement(cust_geom, qn('a:rect'))
        rect.set('l', '0')
        rect.set('t', '0')
        rect.set('r', 'r')
        rect.set('b', 'b')
        # Path must be wrapped in a:pathLst
        path_lst = etree.SubElement(cust_geom, qn('a:pathLst'))
        path_elem = etree.SubElement(path_lst, qn('a:path'))
        path_elem.set('w', str(vw))
        path_elem.set('h', str(vh))

        # Parse SVG path commands
        _parse_svg_path_to_ooxml(svg_path, path_elem)
    except Exception:
        pass
    return shape


def _parse_svg_path_to_ooxml(svg_path, path_elem):
    """Convert SVG path data to OOXML a:path children."""
    import re as _re
    tokens = _re.findall(r'[MmLlHhVvCcSsQqTtAaZz]|[-+]?(?:\d+\.?\d*|\.\d+)', svg_path)
    cx, cy = 0, 0
    i = 0

    def pt(val):
        return str(int(float(val)))

    while i < len(tokens):
        cmd = tokens[i]
        if cmd in ('M', 'm'):
            moveto = etree.SubElement(path_elem, qn('a:moveTo'))
            p = etree.SubElement(moveto, qn('a:pt'))
            p.set('x', pt(tokens[i+1]))
            p.set('y', pt(tokens[i+2]))
            cx, cy = float(tokens[i+1]), float(tokens[i+2])
            i += 3
        elif cmd in ('L', 'l'):
            lineto = etree.SubElement(path_elem, qn('a:lnTo'))
            p = etree.SubElement(lineto, qn('a:pt'))
            if cmd == 'l':
                cx += float(tokens[i+1])
                cy += float(tokens[i+2])
            else:
                cx, cy = float(tokens[i+1]), float(tokens[i+2])
            p.set('x', pt(cx))
            p.set('y', pt(cy))
            i += 3
        elif cmd in ('H', 'h'):
            lineto = etree.SubElement(path_elem, qn('a:lnTo'))
            p = etree.SubElement(lineto, qn('a:pt'))
            if cmd == 'h':
                cx += float(tokens[i+1])
            else:
                cx = float(tokens[i+1])
            p.set('x', pt(cx))
            p.set('y', pt(cy))
            i += 2
        elif cmd in ('V', 'v'):
            lineto = etree.SubElement(path_elem, qn('a:lnTo'))
            p = etree.SubElement(lineto, qn('a:pt'))
            if cmd == 'v':
                cy += float(tokens[i+1])
            else:
                cy = float(tokens[i+1])
            p.set('x', pt(cx))
            p.set('y', pt(cy))
            i += 2
        elif cmd in ('Z', 'z'):
            etree.SubElement(path_elem, qn('a:close'))
            i += 1
        elif cmd in ('C', 'c'):
            # Cubic bezier: 6 params
            pts = []
            for j in range(3):
                bx = float(tokens[i+1+j*2])
                by = float(tokens[i+2+j*2])
                if cmd == 'c':
                    bx += cx
                    by += cy
                pts.append((bx, by))
            cx, cy = pts[-1]
            cubic = etree.SubElement(path_elem, qn('a:cubicBezTo'))
            for bx, by in pts:
                p = etree.SubElement(cubic, qn('a:pt'))
                p.set('x', pt(bx))
                p.set('y', pt(by))
            i += 7
        else:
            # Skip unknown commands
            i += 1


def add_image_element(slide, element, theme, prs=None, base_dir=None):
    bounds = element.get('bounds', [0, 0, 100, 100])
    x, y, w, h = bounds
    src = element.get('src', '')
    if not src:
        return None

    parsed = urlparse(src)
    img_path = None

    if parsed.scheme in ('http', 'https'):
        try:
            import requests
            from io import BytesIO
            resp = requests.get(src, timeout=10)
            if resp.status_code == 200:
                img_path = BytesIO(resp.content)
        except Exception:
            pass
    else:
        p = Path(src)
        if p.exists():
            img_path = str(p)

    is_svg = img_path and isinstance(img_path, str) and img_path.lower().endswith('.svg')
    if is_svg:
        # SVG: inject as p:pic with asvg:svgBlip extension
        rId = _register_svg(slide, img_path)
        shape_id = _next_shape_id(slide)
        display_name = Path(img_path).stem
        pic = _make_pic_svg(px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h),
                            rId, shape_id=shape_id, name=display_name)
        sp_tree = _get_sp_tree(slide)
        sp_tree.append(pic)
        return pic  # SVG: raw lxml element, skip crop/fit/fill/border/shadow/transform
    elif img_path:
        shape = slide.shapes.add_picture(img_path, px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h))
    else:
        # Placeholder — create empty shape (no fill, no border, no text)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h))
        # Remove <p:style> — it adds unwanted theme fill/line references
        style_el = shape._element.find(qn('p:style'))
        if style_el is not None:
            shape._element.remove(style_el)
        shape.fill.background()  # No fill
        shape.line.fill.background()  # No border

    # Crop
    crop = element.get('crop')
    if crop and hasattr(shape, 'crop_left'):
        try:
            shape.crop_left = crop.get('left', 0)
            shape.crop_top = crop.get('top', 0)
            shape.crop_right = crop.get('right', 0)
            shape.crop_bottom = crop.get('bottom', 0)
        except Exception:
            pass

    # Fit mode
    fit = element.get('fit')
    if fit:
        fit_mode = fit.get('mode', 'cover') if isinstance(fit, dict) else fit
        blip_fill = shape._element.find(qn('p:blipFill'))
        if blip_fill is not None:
            if fit_mode == 'tile':
                # Replace stretch with tile
                for stretch in blip_fill.findall(qn('a:stretch')):
                    blip_fill.remove(stretch)
                if blip_fill.find(qn('a:tile')) is None:
                    etree.SubElement(blip_fill, qn('a:tile'))
            elif fit_mode in ('fill', 'cover', 'contain'):
                # Ensure stretch mode (default for fill/cover/contain)
                for tile in blip_fill.findall(qn('a:tile')):
                    blip_fill.remove(tile)
                if blip_fill.find(qn('a:stretch')) is None:
                    etree.SubElement(blip_fill, qn('a:stretch'))

    # Shape clipping (shapeName/adjustments) — modify prstGeom on the picture
    img_shape_name = element.get('shapeName')
    if img_shape_name and img_shape_name != 'rect':
        spPr = shape._element.find(qn('p:spPr'))
        if spPr is not None:
            prst_geom = spPr.find(qn('a:prstGeom'))
            if prst_geom is not None:
                prst_geom.set('prst', img_shape_name)
                # Apply adjustments
                adjustments = element.get('adjustments', [])
                if adjustments:
                    av_lst = prst_geom.find(qn('a:avLst'))
                    if av_lst is None:
                        av_lst = etree.SubElement(prst_geom, qn('a:avLst'))
                    for i, val in enumerate(adjustments):
                        gd = etree.SubElement(av_lst, qn('a:gd'))
                        gd.set('name', f'adj{i + 1}')
                        gd.set('fmla', f'val {val}')

    # Fill and border (for images with borders/outlines)
    fill_spec = element.get('fill')
    if fill_spec:
        apply_fill(shape, fill_spec, theme, prs, base_dir)
    border_spec = element.get('border')
    if border_spec:
        apply_border(shape, border_spec, theme)

    # Shadow
    shadow = element.get('shadow')
    if shadow:
        _inject_shadow(shape, shadow, theme)

    # Opacity
    opacity = element.get('opacity')
    if opacity is not None and opacity < 1:
        apply_opacity(shape, opacity)

    apply_transform(shape, element)
    return shape


def add_icon_element(slide, element, theme, prs=None, base_dir=None):
    """Add icon element. Try SVG from resource, fallback to placeholder."""
    bounds = element.get('bounds', [0, 0, 100, 100])
    x, y, w, h = bounds
    icon_name = element.get('iconName', '')

    svg_path = find_icon_svg(icon_name)

    if svg_path and svg_path.exists():
        shape = slide.shapes.add_picture(str(svg_path), px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h))
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h))
        # Remove <p:style> — it adds unwanted theme fill/line references
        style_el = shape._element.find(qn('p:style'))
        if style_el is not None:
            shape._element.remove(style_el)
        fill_spec = element.get('fill')
        if fill_spec:
            apply_fill(shape, fill_spec, theme, prs, base_dir)
        else:
            shape.fill.background()
        tf = shape.text_frame
        tf.text = f'[{icon_name}]'
        tf.word_wrap = True
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER

    # Border
    border_spec = element.get('border')
    if border_spec:
        apply_border(shape, border_spec, theme)

    # Shadow
    shadow = element.get('shadow')
    if shadow:
        _inject_shadow(shape, shadow, theme)

    # ElementBase: opacity, rotation, flip
    opacity = element.get('opacity')
    if opacity is not None and opacity < 1:
        apply_opacity(shape, opacity)
    apply_transform(shape, element)

    return shape


def add_table_element(slide, element, theme, prs=None, base_dir=None):
    bounds = element.get('bounds', [0, 0, 100, 100])
    x, y, w, h = bounds
    rows_data = element.get('rows', [])
    if not rows_data:
        return None

    num_rows = len(rows_data)
    num_cols = max(len(row) for row in rows_data) if rows_data else 1

    table_width = px_to_emu(w)
    table_height = px_to_emu(h)

    shape = slide.shapes.add_table(num_rows, num_cols, px_to_emu(x), px_to_emu(y), table_width, table_height)
    table = shape.table

    # Apply table style ID if specified
    table_style_id = element.get('tableStyleId') or element.get('tableStyle')
    if table_style_id:
        tblPr = shape._element.find(qn('p:graphicFrame') + '/' + qn('a:graphic') + '/' + qn('a:graphicData') + '/' + qn('a:tbl') + '/' + qn('a:tblPr'))
        if tblPr is None:
            tblPr = shape._element.find('.//' + qn('a:tblPr'))
        if tblPr is not None:
            sid_elem = tblPr.find(qn('a:tableStyleId'))
            if sid_elem is not None:
                sid_elem.text = table_style_id
            else:
                sid_elem = etree.SubElement(tblPr, qn('a:tableStyleId'))
                sid_elem.text = table_style_id

    # Set column widths from percentages
    col_widths = element.get('columnWidths')
    if col_widths and len(col_widths) == num_cols:
        for i, pct in enumerate(col_widths):
            table.columns[i].width = int(table_width * pct)

    # Set row heights from percentages
    row_heights = element.get('rowHeights')
    if row_heights and len(row_heights) == num_rows:
        for i, pct in enumerate(row_heights):
            table.rows[i].height = int(table_height * pct)

    # Table style
    table_style = element.get('style', '$default')
    style_config = None
    if isinstance(table_style, str) and table_style.startswith('$'):
        style_config, _ = resolve_theme_ref(table_style, theme, 'tableStyles')
    elif isinstance(table_style, dict):
        style_config = table_style

    header_fill_color = resolve_color(style_config.get('headerFill'), theme) if style_config else None
    header_text_color = resolve_color(style_config.get('headerColor'), theme) if style_config else None
    body_fill_colors = (style_config.get('bodyFill', []) if style_config else [])
    body_text_color = resolve_color(style_config.get('bodyColor'), theme) if style_config else None
    border_spec = style_config.get('border', {}) if style_config else {}
    border_color = resolve_color(border_spec.get('color'), theme) if isinstance(border_spec, dict) else None
    border_width = border_spec.get('width', 1) if isinstance(border_spec, dict) else 1
    font_family = (style_config.get('fontFamily', 'Arial') if style_config else 'Arial')
    font_size = (style_config.get('fontSize', 14) if style_config else 14)
    header_bold = (style_config.get('headerBold', True) if style_config else True)
    first_col_bold = (style_config.get('firstColumnBold', False) if style_config else False)

    # Additional TableStyleConfig properties
    default_fill_spec = style_config.get('fill') if style_config else None
    default_fill_color = resolve_color(default_fill_spec.get('color'), theme) if isinstance(default_fill_spec, dict) else resolve_color(default_fill_spec, theme) if default_fill_spec else None

    # Border overrides per row type
    header_border_spec = style_config.get('headerBorder') if style_config else None
    body_border_spec = style_config.get('bodyBorder') if style_config else None
    last_row_border_spec = style_config.get('lastRowBorder') if style_config else None

    # First column overrides
    first_col_fill_color = resolve_color(style_config.get('firstColumnFill'), theme) if style_config else None
    first_col_text_color = resolve_color(style_config.get('firstColumnColor'), theme) if style_config else None

    # Track merged cells
    merged = set()  # (row, col) that are consumed by merges

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_data in enumerate(row_data):
            if col_idx >= num_cols or (row_idx, col_idx) in merged:
                continue
            if not isinstance(cell_data, dict):
                cell_data = {'content': {'text': str(cell_data)}}

            # Handle merges
            rs = cell_data.get('rowSpan', 1)
            cs = cell_data.get('colSpan', 1)
            if rs > 1 or cs > 1:
                try:
                    end_row = min(row_idx + rs - 1, num_rows - 1)
                    end_col = min(col_idx + cs - 1, num_cols - 1)
                    merge_cell = table.cell(row_idx, col_idx)
                    merge_cell.merge(table.cell(end_row, end_col))
                    # Mark merged cells
                    for mr in range(row_idx, end_row + 1):
                        for mc in range(col_idx, end_col + 1):
                            if (mr, mc) != (row_idx, col_idx):
                                merged.add((mr, mc))
                except Exception:
                    pass

            cell = table.cell(row_idx, col_idx)

            # Fill
            cell_fill = cell_data.get('fill')
            if cell_fill == 'none':
                # Explicit noFill — add <a:noFill/> to tcPr
                tc_el = cell._tc
                tc_pr = tc_el.find(qn('a:tcPr'))
                if tc_pr is None:
                    tc_pr = etree.SubElement(tc_el, qn('a:tcPr'))
                if tc_pr.find(qn('a:noFill')) is None:
                    # Remove any existing fill elements
                    for fill_tag in ('a:solidFill', 'a:gradFill', 'a:pattFill'):
                        existing = tc_pr.find(qn(fill_tag))
                        if existing is not None:
                            tc_pr.remove(existing)
                    # Insert noFill at the beginning
                    no_fill = etree.Element(qn('a:noFill'))
                    tc_pr.insert(0, no_fill)
            elif cell_fill:
                if isinstance(cell_fill, str):
                    fill_color = resolve_color(cell_fill, theme)
                elif isinstance(cell_fill, dict):
                    fill_color = resolve_color(cell_fill.get('color'), theme)
                else:
                    fill_color = None
                if fill_color:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(fill_color)
                    # Handle hex8 alpha
                    hex_clean = fill_color.lstrip('#')
                    if len(hex_clean) >= 8:
                        alpha_8bit = int(hex_clean[6:8], 16)
                        if alpha_8bit < 255:
                            tc_el = cell._tc
                            tc_pr = tc_el.find(qn('a:tcPr'))
                            if tc_pr is not None:
                                sf = tc_pr.find(qn('a:solidFill'))
                                if sf is not None:
                                    srgb = sf.find(qn('a:srgbClr'))
                                    if srgb is not None:
                                        srgb.set('val', hex_clean[:6])
                                        alpha_elem = srgb.find(qn('a:alpha'))
                                        if alpha_elem is None:
                                            alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                                        alpha_elem.set('val', str(int(alpha_8bit / 255 * 100000)))
            elif col_idx == 0 and first_col_fill_color:
                # First column fill (overrides header/body/default)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(first_col_fill_color)
            elif row_idx == 0 and header_fill_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(header_fill_color)
            elif body_fill_colors:
                fi = (row_idx - 1) % len(body_fill_colors) if row_idx > 0 else 0
                bg = resolve_color(body_fill_colors[fi], theme)
                if bg:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(bg)
            elif default_fill_color:
                # Default fill (lowest priority)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(default_fill_color)

            # Text — use rich text parsing (support both {text:...} and {content:{text:...}} formats)
            content_data = cell_data.get('content', {})
            if isinstance(content_data, str):
                cell_text = content_data
            elif isinstance(content_data, dict):
                cell_text = content_data.get('text', '') or cell_data.get('text', '')
            else:
                cell_text = cell_data.get('text', '')
            if cell_text:
                cell_text = unescape_html(cell_text)
                # Parse rich text for cell
                base_style = {'fontSize': font_size, 'fontFamily': font_family}
                if row_idx == 0 and header_bold:
                    base_style['bold'] = True
                    if header_text_color:
                        base_style['color'] = header_text_color
                elif col_idx == 0 and first_col_bold:
                    base_style['bold'] = True
                    if first_col_text_color:
                        base_style['color'] = first_col_text_color
                if body_text_color and row_idx > 0:
                    base_style['color'] = body_text_color

                segments = parse_rich_text(cell_text, base_style)
                tf = cell.text_frame
                tf.clear()
                current_para = tf.paragraphs[0]
                para_idx = 0
                for seg in segments:
                    if seg.get('is_new_para') and para_idx > 0:
                        current_para = tf.add_paragraph()
                    elif para_idx > 0:
                        pass  # same paragraph
                    para_idx += 1
                    run = current_para.add_run()
                    _set_run_text_with_breaks(run, seg['text'])
                    seg_style = dict(base_style)
                    seg_style.update(seg.get('style', {}))
                    apply_text_style(run, seg_style, theme)
                    run._r.get_or_add_rPr()

            # Per-cell border override
            cell_border = cell_data.get('border')
            if cell_border:
                _apply_cell_border(cell, cell_border, theme)
            else:
                # Row-type specific border (header > body > last row > global)
                row_border_spec = None
                if row_idx == 0 and header_border_spec:
                    row_border_spec = header_border_spec
                elif row_idx == num_rows - 1 and last_row_border_spec:
                    row_border_spec = last_row_border_spec
                elif row_idx > 0 and body_border_spec:
                    row_border_spec = body_border_spec

                if row_border_spec:
                    _apply_cell_border(cell, row_border_spec, theme)
                elif border_color:
                    _apply_cell_border(cell, {'style': 'solid', 'width': border_width, 'color': border_color}, theme)

            # Cell alignment — from content.align [h, v], default [center, middle]
            content_data = cell_data.get('content', {})
            align_spec = content_data.get('align') if isinstance(content_data, dict) else None
            h_align = 'center'
            v_align = 'middle'
            if align_spec and isinstance(align_spec, list) and len(align_spec) >= 2:
                h_align = align_spec[0] or h_align
                v_align = align_spec[1] or v_align
            # Apply vertical anchor
            v_align_map = {'top': 't', 'middle': 'ctr', 'bottom': 'b'}
            anchor_val = v_align_map.get(v_align)
            if anchor_val:
                tc = cell._tc
                tc_pr = tc.find(qn('a:tcPr'))
                if tc_pr is None:
                    tc_pr = etree.SubElement(tc._tc, qn('a:tcPr'))
                tc_pr.set('anchor', anchor_val)
            # Apply horizontal alignment to all paragraphs
            h_align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                           'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY,
                           'distributed': PP_ALIGN.DISTRIBUTE}
            pp_align = h_align_map.get(h_align)
            if pp_align and cell_text:
                for para in cell.text_frame.paragraphs:
                    para.alignment = pp_align

    # Shadow (on the graphicFrame element)
    shadow = element.get('shadow')
    if shadow:
        _inject_shadow(shape, shadow, theme)

    # Column widths
    col_widths = element.get('columnWidths')
    if col_widths and len(col_widths) == num_cols:
        for i, width_pct in enumerate(col_widths):
            table.columns[i].width = int(table_width * width_pct)

    # Row heights
    row_heights = element.get('rowHeights')
    if row_heights and len(row_heights) == num_rows:
        for i, height_pct in enumerate(row_heights):
            table.rows[i].height = int(table_height * height_pct)

    # ElementBase: opacity, rotation, flip
    opacity = element.get('opacity')
    if opacity is not None and opacity < 1:
        apply_opacity(shape, opacity)
    apply_transform(shape, element)

    return shape


def _apply_cell_border(cell, border_spec, theme):
    """Apply border to a table cell via direct XML injection.

    border_spec can be:
    - dict with side keys: {'top': {...}, 'right': {...}, 'bottom': {...}, 'left': {...}}
    - dict applied to all sides: {'style': 'solid', 'width': 1, 'color': '#000'}
    - list of [top, right, bottom, left]
    python-pptx has no cell.border API, so we inject a:tcPr/a:ln* elements.
    """
    side_names = ['left', 'right', 'top', 'bottom', 'tlToBr', 'blToTr']
    side_tags = ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB', 'a:lnTlToBr', 'a:lnBlToTr']

    # Check if border_spec has side keys
    if isinstance(border_spec, dict) and any(k in border_spec for k in side_names):
        borders = [border_spec.get(s, {}) for s in side_names]
    elif isinstance(border_spec, dict):
        borders = [border_spec] * len(side_tags)
    elif isinstance(border_spec, list):
        # Spec format: [top, right, bottom, left] (clockwise)
        # Map to side_tags order: [left, right, top, bottom, tlToBr, blToTr]
        spec_map = {0: 'top', 1: 'right', 2: 'bottom', 3: 'left'}
        borders = []
        for i, side_name in enumerate(side_names):
            if side_name in spec_map.values():
                spec_idx = {v: k for k, v in spec_map.items()}[side_name]
                b = border_spec[spec_idx] if spec_idx < len(border_spec) else None
            else:
                b = None
            borders.append(b if b else {'style': 'none'})
    else:
        return

    # Ensure tcPr exists
    tc = cell._tc
    tc_pr = tc.find(qn('a:tcPr'))
    if tc_pr is None:
        tc_pr = etree.SubElement(tc, qn('a:tcPr'))

    for tag, spec in zip(side_tags, borders):
        if not isinstance(spec, dict):
            continue
        style = spec.get('style', 'solid')
        if style == 'none':
            continue

        # Remove existing border element
        existing = tc_pr.find(qn(tag))
        if existing is not None:
            tc_pr.remove(existing)

        ln = etree.SubElement(tc_pr, qn(tag))
        width = spec.get('width', 1)
        ln.set('w', str(int(width * 12700)))  # pt to EMU

        cmpd = spec.get('cmpd')
        if cmpd:
            ln.set('cmpd', cmpd)

        # Check for noFill inside border
        if spec.get('noFill'):
            etree.SubElement(ln, qn('a:noFill'))
        else:
            color = resolve_color(spec.get('color'), theme)
            if color:
                solid = etree.SubElement(ln, qn('a:solidFill'))
                srgb = etree.SubElement(solid, qn('a:srgbClr'))
                hex_str = color.lstrip('#')
                if len(hex_str) >= 8:
                    srgb.set('val', hex_str[:6])
                    alpha_8bit = int(hex_str[6:8], 16)
                    if alpha_8bit < 255:
                        alpha_pct = round(alpha_8bit / 255 * 100000)
                        alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                        alpha_elem.set('val', str(alpha_pct))
                else:
                    srgb.set('val', hex_str)

        # Border style → prstDash
        prst_dash = spec.get('prstDash')
        if prst_dash:
            dash = etree.SubElement(ln, qn('a:prstDash'))
            dash.set('val', prst_dash)
        elif style in ('dash', 'dot', 'dashDot'):
            dash = etree.SubElement(ln, qn('a:prstDash'))
            dash.set('val', style)


# ---------------------------------------------------------------------------
# Chart XML builders — pure lxml construction of c:chartSpace
# ---------------------------------------------------------------------------

def _build_num_cache(values):
    """Build c:numCache element from a list of numeric values."""
    cache = etree.SubElement(etree.Element('dummy'), qn('c:numCache'))
    etree.SubElement(cache, qn('c:ptCount')).set('val', str(len(values)))
    for i, v in enumerate(values):
        if v is None:
            continue
        pt = etree.SubElement(cache, qn('c:pt'))
        pt.set('idx', str(i))
        etree.SubElement(pt, qn('c:v')).text = str(v)
    return cache


def _build_str_cache(values):
    """Build c:strCache element from a list of string values."""
    cache = etree.SubElement(etree.Element('dummy'), qn('c:strCache'))
    etree.SubElement(cache, qn('c:ptCount')).set('val', str(len(values)))
    for i, v in enumerate(values):
        if v is None:
            continue
        pt = etree.SubElement(cache, qn('c:pt'))
        pt.set('idx', str(i))
        etree.SubElement(pt, qn('c:v')).text = str(v)
    return cache


def _build_num_ref(formula, values):
    """Build c:numRef with formula and numCache."""
    ref = etree.SubElement(etree.Element('dummy'), qn('c:numRef'))
    etree.SubElement(ref, qn('c:f')).text = formula
    ref.append(_build_num_cache(values))
    return ref


def _build_str_ref(formula, values):
    """Build c:strRef with formula and strCache."""
    ref = etree.SubElement(etree.Element('dummy'), qn('c:strRef'))
    etree.SubElement(ref, qn('c:f')).text = formula
    ref.append(_build_str_cache(values))
    return ref


def _build_chart_sp_pr(style, theme):
    """Build c:spPr element for a chart series or axis from a style dict.

    Handles fill (solid/gradient), border (color/width/dash).
    """
    sp_pr = etree.SubElement(etree.Element('dummy'), qn('c:spPr'))

    # Fill
    fill_spec = style.get('fill')
    if fill_spec:
        if isinstance(fill_spec, dict):
            fill_type = fill_spec.get('type', 'solid')
            if fill_type == 'solid':
                color_val = fill_spec.get('color')
                if color_val:
                    solid = etree.SubElement(sp_pr, qn('a:solidFill'))
                    if isinstance(color_val, str) and color_val.startswith('$'):
                        scheme_elem = etree.SubElement(solid, qn('a:schemeClr'))
                        scheme_elem.set('val', color_val[1:])
                        _inject_alpha(scheme_elem, fill_spec.get('alpha'))
                    else:
                        rgb = resolve_color(color_val, theme)
                        if rgb:
                            srgb = etree.SubElement(solid, qn('a:srgbClr'))
                            srgb.set('val', rgb.lstrip('#'))
                            _inject_alpha(srgb, fill_spec.get('alpha'))
            elif fill_type == 'gradient':
                grad = etree.SubElement(sp_pr, qn('a:gradFill'))
                gs_lst = etree.SubElement(grad, qn('a:gsLst'))
                stops = fill_spec.get('stops', [])
                for stop in stops:
                    gs = etree.SubElement(gs_lst, qn('a:gs'))
                    pos = stop.get('position', 0)
                    gs.set('pos', str(int(pos * 100000)))
                    stop_color = stop.get('color')
                    if stop_color:
                        _add_color_element(gs, stop_color, theme)
                lin = etree.SubElement(grad, qn('a:lin'))
                angle = fill_spec.get('angle', 0)
                lin.set('ang', str(int(angle * 60000)))
                lin.set('scaled', '1')
        elif isinstance(fill_spec, str):
            _add_color_element(sp_pr, fill_spec, theme)

    # Border/Line
    border_spec = style.get('border')
    if border_spec:
        ln = etree.SubElement(sp_pr, qn('a:ln'))
        if isinstance(border_spec, dict):
            if border_spec.get('style') == 'none':
                etree.SubElement(ln, qn('a:noFill'))
            else:
                w = border_spec.get('width')
                if w:
                    ln.set('w', str(int(w * 12700)))
                c = border_spec.get('color')
                if c:
                    _add_color_element(ln, c, theme)
                s = border_spec.get('style')
                if s and s != 'none':
                    etree.SubElement(ln, qn('a:prstDash')).set('val', s)
        elif isinstance(border_spec, str):
            _add_color_element(ln, border_spec, theme)

    return sp_pr


def _inject_alpha(color_elem, alpha_val):
    """Inject a:alpha element into a color element if alpha < 1."""
    if alpha_val is not None and alpha_val < 1:
        alpha_pct = int(alpha_val * 100000)
        alpha = etree.SubElement(color_elem, qn('a:alpha'))
        alpha.set('val', str(alpha_pct))


def _add_color_element(parent, color_value, theme):
    """Add a:solidFill with a:schemeClr or a:srgbClr to parent element.

    color_value can be:
    - '$key' → a:schemeClr (theme reference)
    - '#RRGGBB' → a:srgbClr (direct hex)
    """
    if not color_value:
        return
    solid = etree.SubElement(parent, qn('a:solidFill'))
    if isinstance(color_value, str) and color_value.startswith('$'):
        scheme_name = color_value[1:]
        etree.SubElement(solid, qn('a:schemeClr')).set('val', scheme_name)
    else:
        rgb = resolve_color(color_value, theme)
        if rgb:
            etree.SubElement(solid, qn('a:srgbClr')).set('val', rgb.lstrip('#'))


def _build_chart_dlbls(dl_spec, theme):
    """Build c:dLbls element from a DataLabelConfig dict."""
    if not dl_spec or not isinstance(dl_spec, dict):
        return None

    dlbls = etree.SubElement(etree.Element('dummy'), qn('c:dLbls'))

    show = dl_spec.get('show', False)
    if not show:
        etree.SubElement(dlbls, qn('c:delete')).set('val', '1')
        return dlbls

    etree.SubElement(dlbls, qn('c:showLegendKey')).set('val', '0')

    content = dl_spec.get('content', 'value')
    etree.SubElement(dlbls, qn('c:showVal')).set('val', '1' if content == 'value' else '0')
    etree.SubElement(dlbls, qn('c:showCatName')).set('val', '1' if content == 'category' else '0')
    etree.SubElement(dlbls, qn('c:showSerName')).set('val', '1' if content == 'name' else '0')
    etree.SubElement(dlbls, qn('c:showPercent')).set('val', '1' if content == 'percentage' else '0')
    etree.SubElement(dlbls, qn('c:showBubbleSize')).set('val', '0')

    nf = dl_spec.get('numberFormat')
    if nf:
        num_fmt = etree.SubElement(dlbls, qn('c:numFmt'))
        num_fmt.set('formatCode', nf)
        num_fmt.set('sourceLinked', '0')

    dl_color = dl_spec.get('color')
    dl_fs = dl_spec.get('fontSize')
    if dl_color or dl_fs:
        tx_pr = etree.SubElement(dlbls, qn('c:txPr'))
        etree.SubElement(tx_pr, qn('a:bodyPr'))
        etree.SubElement(tx_pr, qn('a:lstStyle'))
        p = etree.SubElement(tx_pr, qn('a:p'))
        p_pr = etree.SubElement(p, qn('a:pPr'))
        def_rpr = etree.SubElement(p_pr, qn('a:defRPr'))
        if dl_fs:
            def_rpr.set('sz', str(int(dl_fs * 100)))
        if dl_color:
            _add_color_element(def_rpr, dl_color, theme)
        etree.SubElement(p, qn('a:endParaRPr')).set('lang', 'en-US')

    return dlbls


def _build_chart_marker(marker_spec, theme):
    """Build c:marker element from a MarkerConfig."""
    if marker_spec is None:
        return None

    marker = etree.SubElement(etree.Element('dummy'), qn('c:marker'))

    if marker_spec is False:
        etree.SubElement(marker, qn('c:symbol')).set('val', 'none')
        return marker

    if isinstance(marker_spec, dict):
        shape = marker_spec.get('shape', 'circle')
        etree.SubElement(marker, qn('c:symbol')).set('val', shape)

        size = marker_spec.get('size')
        if size:
            etree.SubElement(marker, qn('c:size')).set('val', str(size))

        m_fill = marker_spec.get('fill')
        m_border = marker_spec.get('border')
        if m_fill or m_border:
            sp_pr = etree.SubElement(marker, qn('c:spPr'))
            if m_fill:
                _add_color_element(sp_pr, m_fill, theme)
            if m_border and isinstance(m_border, dict):
                ln = etree.SubElement(sp_pr, qn('a:ln'))
                b_width = m_border.get('width')
                if b_width:
                    ln.set('w', str(int(b_width * 12700)))
                b_color = m_border.get('color')
                if b_color:
                    _add_color_element(ln, b_color, theme)

    return marker


def _build_chart_ser(idx, name, categories, values, style=None, per_point_colors=None,
                     theme=None, chart_type='bar'):
    """Build a complete c:ser element.

    Args:
        idx: Series index
        name: Series name string
        categories: List of category labels (for cat/val charts)
        values: List of numeric values
        style: SeriesStyleConfig dict
        per_point_colors: List of resolved hex colors for data points
        theme: Theme dict
        chart_type: 'bar', 'line', 'area', 'scatter', etc.
    """
    ser = etree.SubElement(etree.Element('dummy'), qn('c:ser'))
    etree.SubElement(ser, qn('c:idx')).set('val', str(idx))
    etree.SubElement(ser, qn('c:order')).set('val', str(idx))

    # Series name (c:tx / c:strRef)
    if name:
        tx = etree.SubElement(ser, qn('c:tx'))
        str_ref = etree.SubElement(tx, qn('c:strRef'))
        etree.SubElement(str_ref, qn('c:f')).text = f'Sheet1!$B${idx + 1}'
        str_cache = etree.SubElement(str_ref, qn('c:strCache'))
        etree.SubElement(str_cache, qn('c:ptCount')).set('val', '1')
        pt = etree.SubElement(str_cache, qn('c:pt'))
        pt.set('idx', '0')
        etree.SubElement(pt, qn('c:v')).text = str(name)

    # SpPr (series styling)
    if style:
        sp_pr = _build_chart_sp_pr(style, theme)
        if len(sp_pr):
            ser.append(sp_pr)

    # Category and value references
    use_xy = chart_type in ('scatter', 'bubble')

    if use_xy:
        # X values (c:xVal / c:numRef)
        if categories:
            x_val = etree.SubElement(ser, qn('c:xVal'))
            x_ref = etree.SubElement(x_val, qn('c:numRef'))
            etree.SubElement(x_ref, qn('c:f')).text = f'Sheet1!$A$2:$A${len(categories)+1}'
            x_ref.append(_build_num_cache(categories))

        # Y values (c:yVal / c:numRef)
        y_val = etree.SubElement(ser, qn('c:yVal'))
        y_ref = etree.SubElement(y_val, qn('c:numRef'))
        etree.SubElement(y_ref, qn('c:f')).text = f'Sheet1!$B$2:$B${len(values)+1}'
        y_ref.append(_build_num_cache(values))
    else:
        # Category (c:cat / c:strRef)
        if categories:
            cat = etree.SubElement(ser, qn('c:cat'))
            cat_ref = etree.SubElement(cat, qn('c:strRef'))
            etree.SubElement(cat_ref, qn('c:f')).text = f'Sheet1!$A$2:$A${len(categories)+1}'
            cat_ref.append(_build_str_cache([str(c) for c in categories]))

        # Values (c:val / c:numRef)
        val = etree.SubElement(ser, qn('c:val'))
        val_ref = etree.SubElement(val, qn('c:numRef'))
        etree.SubElement(val_ref, qn('c:f')).text = f'Sheet1!$B$2:$B${len(values)+1}'
        val_ref.append(_build_num_cache(values))

    # Per-data-point colors (c:dPt)
    if per_point_colors:
        for i, color in enumerate(per_point_colors):
            if color:
                dpt = etree.SubElement(ser, qn('c:dpt'))
                etree.SubElement(dpt, qn('c:idx')).set('val', str(i))
                etree.SubElement(dpt, qn('c:bubble3D')).set('val', '0')
                dpt_sp = etree.SubElement(dpt, qn('c:spPr'))
                _add_color_element(dpt_sp, color, theme)
                dpt_ln = etree.SubElement(dpt_sp, qn('a:ln'))
                etree.SubElement(dpt_ln, qn('a:noFill'))

    # Data labels (c:dLbls) — series-level
    if style:
        series_dl = style.get('dataLabels')
        if series_dl:
            dlbls = _build_chart_dlbls(series_dl, theme)
            if dlbls is not None:
                ser.append(dlbls)

    # Smooth (for line/area)
    if style and 'smooth' in style:
        etree.SubElement(ser, qn('c:smooth')).set('val', '1' if style['smooth'] else '0')

    # Marker (for line charts)
    if style:
        marker_spec = style.get('marker')
        if marker_spec is not None:
            marker = _build_chart_marker(marker_spec, theme)
            if marker is not None:
                ser.append(marker)

    return ser


def _build_chart_legend(legend_spec, theme, font_family=None):
    """Build c:legend element from a LegendConfig."""
    if legend_spec is None:
        return None

    legend = etree.SubElement(etree.Element('dummy'), qn('c:legend'))

    # Position
    if isinstance(legend_spec, dict):
        pos = legend_spec.get('position', 'bottom')
    else:
        pos = 'bottom'

    pos_map = {
        'top': 't', 'bottom': 'b', 'left': 'l', 'right': 'r',
        'topRight': 'tr',
    }
    etree.SubElement(legend, qn('c:legendPos')).set('val', pos_map.get(pos, 'b'))

    # Overlay
    etree.SubElement(legend, qn('c:overlay')).set('val', '0')

    # Text properties
    legend_color = legend_spec.get('color') if isinstance(legend_spec, dict) else None
    legend_fs = legend_spec.get('fontSize') if isinstance(legend_spec, dict) else None
    if legend_color or legend_fs or font_family:
        tx_pr = etree.SubElement(legend, qn('c:txPr'))
        etree.SubElement(tx_pr, qn('a:bodyPr'))
        etree.SubElement(tx_pr, qn('a:lstStyle'))
        p = etree.SubElement(tx_pr, qn('a:p'))
        p_pr = etree.SubElement(p, qn('a:pPr'))
        def_rpr = etree.SubElement(p_pr, qn('a:defRPr'))
        if legend_fs:
            def_rpr.set('sz', str(int(legend_fs * 100)))
        if font_family:
            def_rpr.set('typeface', font_family)
        if legend_color:
            _add_color_element(def_rpr, legend_color, theme)
        etree.SubElement(p, qn('a:endParaRPr')).set('lang', 'en-US')

    return legend


def _build_chart_title(title_spec, theme, font_family=None):
    """Build c:title element from a title spec (string or dict)."""
    title = etree.SubElement(etree.Element('dummy'), qn('c:title'))

    tx = etree.SubElement(title, qn('c:tx'))
    rich = etree.SubElement(tx, qn('c:rich'))
    etree.SubElement(rich, qn('a:bodyPr'))
    etree.SubElement(rich, qn('a:lstStyle'))
    p = etree.SubElement(rich, qn('a:p'))

    r = etree.SubElement(p, qn('a:r'))
    rpr = etree.SubElement(r, qn('a:rPr'))
    rpr.set('lang', 'en-US')
    rpr.set('dirty', '0')

    if isinstance(title_spec, dict):
        text = title_spec.get('text', '')
        title_color = title_spec.get('color')
        title_fs = title_spec.get('fontSize')
    else:
        text = str(title_spec)
        title_color = None
        title_fs = None

    if title_fs:
        rpr.set('sz', str(int(title_fs * 100)))
    if font_family:
        rpr.set('typeface', font_family)
    if title_color:
        _add_color_element(rpr, title_color, theme)

    etree.SubElement(r, qn('a:t')).text = text

    etree.SubElement(title, qn('c:overlay')).set('val', '0')

    return title


def _build_cat_axis(axis_id, cross_ax_id, position, config, theme, font_family=None):
    """Build c:catAx element."""
    ax = etree.SubElement(etree.Element('dummy'), qn('c:catAx'))
    etree.SubElement(ax, qn('c:axId')).set('val', str(axis_id))

    scaling = etree.SubElement(ax, qn('c:scaling'))
    etree.SubElement(scaling, qn('c:orientation')).set('val', 'minMax')

    if config.get('show') is False:
        etree.SubElement(ax, qn('c:delete')).set('val', '1')
    else:
        etree.SubElement(ax, qn('c:delete')).set('val', '0')

    etree.SubElement(ax, qn('c:axPos')).set('val', position)

    # Gridlines
    gl = config.get('gridLine')
    if gl is True or isinstance(gl, dict):
        major_gl = etree.SubElement(ax, qn('c:majorGridlines'))
        if isinstance(gl, dict):
            sp_pr = etree.SubElement(major_gl, qn('c:spPr'))
            ln = etree.SubElement(sp_pr, qn('a:ln'))
            gl_width = gl.get('width')
            if gl_width:
                ln.set('w', str(int(gl_width * 12700)))
            gl_color = gl.get('color')
            if gl_color:
                _add_color_element(ln, gl_color, theme)
            gl_style = gl.get('style')
            if gl_style:
                etree.SubElement(ln, qn('a:prstDash')).set('val', gl_style)

    etree.SubElement(ax, qn('c:numFmt')).set('formatCode', 'General')
    etree.SubElement(ax, qn('c:numFmt')).set('sourceLinked', '0')
    etree.SubElement(ax, qn('c:majorTickMark')).set('val', 'out')
    etree.SubElement(ax, qn('c:minorTickMark')).set('val', 'none')
    etree.SubElement(ax, qn('c:tickLblPos')).set('val', 'nextTo')

    # Axis line (c:spPr)
    axis_line = config.get('axisLine')
    if axis_line is not None:
        sp_pr = etree.SubElement(ax, qn('c:spPr'))
        if axis_line is False:
            etree.SubElement(sp_pr, qn('a:noFill'))
        elif isinstance(axis_line, dict):
            ln = etree.SubElement(sp_pr, qn('a:ln'))
            al_width = axis_line.get('width')
            if al_width:
                ln.set('w', str(int(al_width * 12700)))
            al_color = axis_line.get('color')
            if al_color:
                _add_color_element(ln, al_color, theme)
            al_style = axis_line.get('style')
            if al_style:
                etree.SubElement(ln, qn('a:prstDash')).set('val', al_style)
            arrow = axis_line.get('arrow')
            if arrow and isinstance(arrow, list):
                if len(arrow) >= 1 and arrow[0]:
                    head = etree.SubElement(ln, qn('a:headEnd'))
                    head.set('type', arrow[0])
                if len(arrow) >= 2 and arrow[1]:
                    tail = etree.SubElement(ln, qn('a:tailEnd'))
                    tail.set('type', arrow[1])

    # Cross axis
    etree.SubElement(ax, qn('c:crossAx')).set('val', str(cross_ax_id))
    etree.SubElement(ax, qn('c:crosses')).set('val', 'autoZero')

    # Tick label font
    label = config.get('label')
    if label is False:
        etree.SubElement(ax, qn('c:txPr'))
    elif isinstance(label, dict) or font_family:
        tx_pr = etree.SubElement(ax, qn('c:txPr'))
        etree.SubElement(tx_pr, qn('a:bodyPr'))
        etree.SubElement(tx_pr, qn('a:lstStyle'))
        p = etree.SubElement(tx_pr, qn('a:p'))
        p_pr = etree.SubElement(p, qn('a:pPr'))
        def_rpr = etree.SubElement(p_pr, qn('a:defRPr'))
        if isinstance(label, dict):
            label_fs = label.get('fontSize')
            if label_fs:
                def_rpr.set('sz', str(int(label_fs * 100)))
            label_color = label.get('color')
            if label_color:
                _add_color_element(def_rpr, label_color, theme)
        if font_family:
            def_rpr.set('typeface', font_family)
        etree.SubElement(p, qn('a:endParaRPr')).set('lang', 'en-US')

    # Title
    title_spec = config.get('title')
    if title_spec:
        title_elem = _build_axis_title(title_spec, theme, font_family)
        if title_elem is not None:
            ax.append(title_elem)

    etree.SubElement(ax, qn('c:crossBetween')).set('val', 'between')

    return ax


def _build_val_axis(axis_id, cross_ax_id, position, config, theme, font_family=None):
    """Build c:valAx element."""
    ax = etree.SubElement(etree.Element('dummy'), qn('c:valAx'))
    etree.SubElement(ax, qn('c:axId')).set('val', str(axis_id))

    scaling = etree.SubElement(ax, qn('c:scaling'))
    etree.SubElement(scaling, qn('c:orientation')).set('val', 'minMax')

    # Min/Max
    if 'min' in config:
        etree.SubElement(scaling, qn('c:min')).set('val', str(config['min']))
    if 'max' in config:
        etree.SubElement(scaling, qn('c:max')).set('val', str(config['max']))

    if config.get('show') is False:
        etree.SubElement(ax, qn('c:delete')).set('val', '1')
    else:
        etree.SubElement(ax, qn('c:delete')).set('val', '0')

    etree.SubElement(ax, qn('c:axPos')).set('val', position)

    # Gridlines
    gl = config.get('gridLine')
    if gl is True or isinstance(gl, dict):
        major_gl = etree.SubElement(ax, qn('c:majorGridlines'))
        if isinstance(gl, dict):
            sp_pr = etree.SubElement(major_gl, qn('c:spPr'))
            ln = etree.SubElement(sp_pr, qn('a:ln'))
            gl_width = gl.get('width')
            if gl_width:
                ln.set('w', str(int(gl_width * 12700)))
            gl_color = gl.get('color')
            if gl_color:
                _add_color_element(ln, gl_color, theme)
            gl_style = gl.get('style')
            if gl_style:
                etree.SubElement(ln, qn('a:prstDash')).set('val', gl_style)

    nf = config.get('numberFormat', 'General')
    num_fmt = etree.SubElement(ax, qn('c:numFmt'))
    num_fmt.set('formatCode', nf)
    num_fmt.set('sourceLinked', '0')

    etree.SubElement(ax, qn('c:majorTickMark')).set('val', 'out')
    etree.SubElement(ax, qn('c:minorTickMark')).set('val', 'none')
    etree.SubElement(ax, qn('c:tickLblPos')).set('val', 'nextTo')

    # Axis line (c:spPr)
    axis_line = config.get('axisLine')
    if axis_line is not None:
        sp_pr = etree.SubElement(ax, qn('c:spPr'))
        if axis_line is False:
            etree.SubElement(sp_pr, qn('a:noFill'))
        elif isinstance(axis_line, dict):
            ln = etree.SubElement(sp_pr, qn('a:ln'))
            al_width = axis_line.get('width')
            if al_width:
                ln.set('w', str(int(al_width * 12700)))
            al_color = axis_line.get('color')
            if al_color:
                _add_color_element(ln, al_color, theme)
            al_style = axis_line.get('style')
            if al_style:
                etree.SubElement(ln, qn('a:prstDash')).set('val', al_style)
            arrow = axis_line.get('arrow')
            if arrow and isinstance(arrow, list):
                if len(arrow) >= 1 and arrow[0]:
                    head = etree.SubElement(ln, qn('a:headEnd'))
                    head.set('type', arrow[0])
                if len(arrow) >= 2 and arrow[1]:
                    tail = etree.SubElement(ln, qn('a:tailEnd'))
                    tail.set('type', arrow[1])

    # Cross axis
    etree.SubElement(ax, qn('c:crossAx')).set('val', str(cross_ax_id))
    etree.SubElement(ax, qn('c:crosses')).set('val', 'autoZero')

    # Tick label font
    label = config.get('label')
    if label is False:
        etree.SubElement(ax, qn('c:txPr'))
    elif isinstance(label, dict) or font_family:
        tx_pr = etree.SubElement(ax, qn('c:txPr'))
        etree.SubElement(tx_pr, qn('a:bodyPr'))
        etree.SubElement(tx_pr, qn('a:lstStyle'))
        p = etree.SubElement(tx_pr, qn('a:p'))
        p_pr = etree.SubElement(p, qn('a:pPr'))
        def_rpr = etree.SubElement(p_pr, qn('a:defRPr'))
        if isinstance(label, dict):
            label_fs = label.get('fontSize')
            if label_fs:
                def_rpr.set('sz', str(int(label_fs * 100)))
            label_color = label.get('color')
            if label_color:
                _add_color_element(def_rpr, label_color, theme)
        if font_family:
            def_rpr.set('typeface', font_family)
        etree.SubElement(p, qn('a:endParaRPr')).set('lang', 'en-US')

    # Title
    title_spec = config.get('title')
    if title_spec:
        title_elem = _build_axis_title(title_spec, theme, font_family)
        if title_elem is not None:
            ax.append(title_elem)

    return ax


def _build_axis_title(title_spec, theme, font_family=None):
    """Build c:title element for an axis title."""
    title = etree.SubElement(etree.Element('dummy'), qn('c:title'))

    tx = etree.SubElement(title, qn('c:tx'))
    rich = etree.SubElement(tx, qn('c:rich'))
    etree.SubElement(rich, qn('a:bodyPr'))
    etree.SubElement(rich, qn('a:lstStyle'))
    p = etree.SubElement(rich, qn('a:p'))

    r = etree.SubElement(p, qn('a:r'))
    rpr = etree.SubElement(r, qn('a:rPr'))
    rpr.set('lang', 'en-US')
    rpr.set('dirty', '0')

    if isinstance(title_spec, dict):
        text = title_spec.get('text', '')
        title_color = title_spec.get('color')
        title_fs = title_spec.get('fontSize')
    else:
        text = str(title_spec)
        title_color = None
        title_fs = None

    if title_fs:
        rpr.set('sz', str(int(title_fs * 100)))
    if font_family:
        rpr.set('typeface', font_family)
    if title_color:
        _add_color_element(rpr, title_color, theme)

    etree.SubElement(r, qn('a:t')).text = text

    etree.SubElement(title, qn('c:overlay')).set('val', '0')

    return title


_SENTINEL = object()


def _build_chart_xml(chart_type, categories, y_field_names, series_data,
                     options=None, colors=None, series_style=None,
                     data_labels=None, legend_spec=_SENTINEL, title_spec=None,
                     x_axis=None, y_axis=None, secondary_axis=None,
                     theme=None, font_family=None):
    """Build complete c:chartSpace XML bytes.

    Args:
        chart_type: 'bar', 'line', 'area', 'pie', 'scatter', 'bubble', 'radar', 'combo'
        categories: List of category labels
        y_field_names: List of series names
        series_data: List of (name, x_values, y_values) per series
        options: ChartOptions dict
        colors: List of color strings
        series_style: Dict of series style configs
        data_labels: DataLabelConfig dict
        legend_spec: LegendConfig or None or _SENTINEL (not specified)
        title_spec: Title spec
        x_axis: AxisConfig for category axis
        y_axis: AxisConfig for value axis
        secondary_axis: AxisConfig for secondary axis
        theme: Theme dict
        font_family: Global font family string
    """
    if options is None:
        options = {}
    if series_style is None:
        series_style = {}
    if theme is None:
        theme = {}
    if x_axis is None:
        x_axis = {}
    if y_axis is None:
        y_axis = {}

    is_doughnut = chart_type == 'pie' and 'innerRadius' in options

    # --- c:chartSpace root ---
    cs = etree.Element(qn('c:chartSpace'), nsmap={
        'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    })

    # ChartSpace-level metadata
    etree.SubElement(cs, qn('c:date1904')).set('val', '0')
    etree.SubElement(cs, qn('c:lang')).set('val', 'zh-CN')
    etree.SubElement(cs, qn('c:roundedCorners')).set('val', '0')

    chart = etree.SubElement(cs, qn('c:chart'))

    # --- Title / autoTitleDeleted ---
    if title_spec:
        chart.append(_build_chart_title(title_spec, theme, font_family))
    else:
        etree.SubElement(chart, qn('c:autoTitleDeleted')).set('val', '1')

    # --- c:plotArea ---
    plot_area = etree.SubElement(chart, qn('c:plotArea'))
    etree.SubElement(plot_area, qn('c:layout'))

    # --- Build series ---
    resolved_series = []
    for i, y_field in enumerate(y_field_names):
        if i >= len(series_data):
            break
        name, x_vals, y_vals = series_data[i]

        # Style resolution
        style = series_style.get(y_field, {})
        if not style:
            style = series_style.get('*', {})
        if not style and i < len(series_data):
            style = series_data[i][3] if len(series_data[i]) > 3 else {}

        # Series name override from seriesStyle
        if style and 'name' in style:
            name = style['name']

        # Per-point colors
        per_point_colors = None
        if colors:
            if chart_type in ('pie', 'doughnut') or is_doughnut:
                per_point_colors = colors
            else:
                if i < len(colors):
                    style = dict(style) if style else {}
                    style.setdefault('fill', {})
                    if isinstance(style.get('fill'), dict):
                        style['fill']['color'] = colors[i]
                    else:
                        style['fill'] = {'type': 'solid', 'color': colors[i]}

        ser = _build_chart_ser(
            i, name, x_vals, y_vals,
            style=style, per_point_colors=per_point_colors,
            theme=theme, chart_type=chart_type,
        )
        resolved_series.append((ser, style))

    # --- Build chart type elements ---
    oxml_chart_tag = {
        'bar': 'c:barChart', 'line': 'c:lineChart', 'area': 'c:areaChart',
        'pie': 'c:doughnutChart' if is_doughnut else 'c:pieChart',
        'scatter': 'c:scatterChart', 'bubble': 'c:bubbleChart',
        'radar': 'c:radarChart',
    }

    use_xy = chart_type in ('scatter', 'bubble')
    _CAT_AXIS_ID = 111111111
    _VAL_AXIS_ID = 222222222
    _SEC_AXIS_ID = 333333333
    needs_axes = chart_type in ('bar', 'line', 'area', 'combo') and not is_doughnut

    if chart_type == 'combo':
        # Group series by type
        type_groups = {}
        for i, (ser, style) in enumerate(resolved_series):
            s_type = style.get('type', 'bar')
            axis = style.get('axis', 'primary')
            key = (s_type, axis)
            if key not in type_groups:
                type_groups[key] = []
            type_groups[key].append(ser)

        combo_tags = {'bar': 'c:barChart', 'line': 'c:lineChart', 'area': 'c:areaChart'}
        for (s_type, axis), sers in type_groups.items():
            tag = combo_tags.get(s_type, 'c:barChart')
            ct = etree.SubElement(plot_area, qn(tag))

            grouping = options.get('stacked')
            if grouping == '100%':
                etree.SubElement(ct, qn('c:grouping')).set('val', 'percentStacked')
            elif grouping:
                etree.SubElement(ct, qn('c:grouping')).set('val', 'stacked')
            else:
                etree.SubElement(ct, qn('c:grouping')).set('val', 'standard')

            if s_type == 'bar':
                bar_dir = options.get('direction', 'vertical')
                etree.SubElement(ct, qn('c:barDir')).set('val', 'bar' if bar_dir == 'horizontal' else 'col')

            for ser in sers:
                ct.append(ser)

            # VaryColors for pie/doughnut
            if s_type in ('pie', 'doughnut'):
                etree.SubElement(ct, qn('c:varyColors')).set('val', '1')

            # AxId references
            etree.SubElement(ct, qn('c:axId')).set('val', str(_CAT_AXIS_ID))
            if axis == 'secondary':
                etree.SubElement(ct, qn('c:axId')).set('val', str(_SEC_AXIS_ID))
            else:
                etree.SubElement(ct, qn('c:axId')).set('val', str(_VAL_AXIS_ID))

        # Overlap for stacked bar
        if options.get('stacked'):
            for ct in plot_area:
                if ct.tag == qn('c:barChart'):
                    grouping = ct.find(qn('c:grouping'))
                    if grouping is not None and grouping.get('val') in ('stacked', 'percentStacked'):
                        overlap_elem = ct.find(qn('c:overlap'))
                        if overlap_elem is None:
                            overlap_elem = etree.SubElement(ct, qn('c:overlap'))
                        overlap_elem.set('val', '100')
                    break

    else:
        # Single chart type
        tag = oxml_chart_tag.get(chart_type, 'c:barChart')
        ct = etree.SubElement(plot_area, qn(tag))

        if chart_type in ('bar', 'line', 'area'):
            grouping = options.get('stacked')
            if grouping == '100%':
                etree.SubElement(ct, qn('c:grouping')).set('val', 'percentStacked')
            elif grouping:
                etree.SubElement(ct, qn('c:grouping')).set('val', 'stacked')
            else:
                etree.SubElement(ct, qn('c:grouping')).set('val', 'standard')

        if chart_type == 'bar':
            bar_dir = options.get('direction', 'vertical')
            etree.SubElement(ct, qn('c:barDir')).set('val', 'bar' if bar_dir == 'horizontal' else 'col')

        if is_doughnut:
            hole_pct = int(options['innerRadius'] * 100)
            etree.SubElement(ct, qn('c:holeSize')).set('val', str(max(10, min(90, hole_pct))))

        if (chart_type in ('pie',) or is_doughnut) and 'startAngle' in options:
            etree.SubElement(ct, qn('c:firstSliceAng')).set('val', str(int(options['startAngle'])))

        if chart_type == 'scatter':
            null_handling = options.get('nullHandling')
            if null_handling == 'connect':
                etree.SubElement(ct, qn('c:scatterStyle')).set('val', 'line')
            else:
                etree.SubElement(ct, qn('c:scatterStyle')).set('val', 'marker')

        if chart_type in ('pie',) or is_doughnut:
            etree.SubElement(ct, qn('c:varyColors')).set('val', '1')

        # Bar width
        bar_width = options.get('barWidth')
        if bar_width and chart_type == 'bar':
            bar_pct = max(0, min(100, int(bar_width * 100)))
            etree.SubElement(ct, qn('c:barSz')).set('val', str(bar_pct))

        for ser, _ in resolved_series:
            ct.append(ser)

        # AxId references (not for pie/doughnut/radar)
        if needs_axes:
            etree.SubElement(ct, qn('c:axId')).set('val', str(_CAT_AXIS_ID))
            etree.SubElement(ct, qn('c:axId')).set('val', str(_VAL_AXIS_ID))
        elif use_xy:
            etree.SubElement(ct, qn('c:axId')).set('val', str(_CAT_AXIS_ID))
            etree.SubElement(ct, qn('c:axId')).set('val', str(_VAL_AXIS_ID))

        # Overlap for stacked bar
        if chart_type == 'bar' and options.get('stacked'):
            overlap_elem = ct.find(qn('c:overlap'))
            if overlap_elem is None:
                overlap_elem = etree.SubElement(ct, qn('c:overlap'))
            overlap_elem.set('val', '100')

    # --- Axes ---
    if needs_axes:
        plot_area.append(_build_cat_axis(_CAT_AXIS_ID, _VAL_AXIS_ID, 'b', x_axis, theme, font_family))
        plot_area.append(_build_val_axis(_VAL_AXIS_ID, _CAT_AXIS_ID, 'l', y_axis, theme, font_family))

        # Secondary axis (combo charts)
        if secondary_axis:
            plot_area.append(_build_val_axis(_SEC_AXIS_ID, _CAT_AXIS_ID, 'r', secondary_axis, theme, font_family))

    elif use_xy:
        # Scatter/bubble: two value axes
        x_config = dict(x_axis)
        x_config.setdefault('numberFormat', 'General')
        plot_area.append(_build_val_axis(_CAT_AXIS_ID, _VAL_AXIS_ID, 'b', x_config, theme, font_family))
        plot_area.append(_build_val_axis(_VAL_AXIS_ID, _CAT_AXIS_ID, 'l', y_axis, theme, font_family))

    # --- Legend ---
    if legend_spec is _SENTINEL:
        legend_spec = None  # Not specified → no legend
    if legend_spec is not None:
        show = True
        if isinstance(legend_spec, bool):
            show = legend_spec
        elif isinstance(legend_spec, dict):
            show = legend_spec.get('show', True)
        if show:
            legend = _build_chart_legend(legend_spec, theme, font_family)
            if legend is not None:
                chart.append(legend)

    # --- Display blanks ---
    null_handling = options.get('nullHandling')
    blanks_map = {'connect': 'span', 'gap': 'gap', 'zero': 'zero'}
    etree.SubElement(chart, qn('c:dispBlanksAs')).set('val', blanks_map.get(null_handling, 'gap'))

    etree.SubElement(chart, qn('c:plotVisOnly')).set('val', '1')
    etree.SubElement(chart, qn('c:showDLblsOverMax')).set('val', '0')

    # --- Chart-level data labels ---
    if data_labels and isinstance(data_labels, dict):
        dlbls = _build_chart_dlbls(data_labels, theme)
        if dlbls is not None:
            # Attach to the first plot area child (chart type element)
            for ct in plot_area:
                if ct.tag not in (qn('c:layout'), qn('c:catAx'), qn('c:valAx'), qn('c:serAx')):
                    ct.append(dlbls)
                    break

    # Serialize
    return etree.tostring(cs, xml_declaration=True, encoding='UTF-8', standalone=True)


def add_chart_element(slide, element, theme, prs=None, base_dir=None):
    """Add a chart element to the slide using python-pptx for OPC + lxml for XML."""
    bounds = element.get('bounds', [0, 0, 100, 100])
    x, y, w, h = bounds

    from pptx.chart.data import ChartData, BubbleChartData
    from pptx.enum.chart import XL_CHART_TYPE

    # --- Resolve chart type ---
    chart_type = element.get('type', element.get('chartType', 'bar'))
    options = element.get('options', {})
    is_doughnut = chart_type == 'pie' and 'innerRadius' in options
    # --- Resolve XL chart type for OPC skeleton ---
    chart_type_map = {
        'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE,
        'area': XL_CHART_TYPE.AREA,
        'pie': XL_CHART_TYPE.PIE,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
        'scatter': XL_CHART_TYPE.XY_SCATTER,
        'radar': XL_CHART_TYPE.RADAR,
        'bubble': XL_CHART_TYPE.BUBBLE,
        'combo': XL_CHART_TYPE.COLUMN_CLUSTERED,
    }
    if is_doughnut:
        xl_type = XL_CHART_TYPE.DOUGHNUT
    else:
        xl_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    bar_dir = options.get('direction') or element.get('barDirection', '')
    if chart_type == 'bar' and bar_dir == 'horizontal':
        xl_type = XL_CHART_TYPE.BAR_CLUSTERED

    # --- Parse data for Excel blob ---
    data = element.get('data', [])
    if not data:
        return None

    _first = data[0] if data else {}
    _is_series_oriented = isinstance(_first, dict) and 'x' in _first and 'y' in _first
    y_field_names = []
    size_field = element.get('size')

    if chart_type == 'bubble':
        chart_data = BubbleChartData()
        if _is_series_oriented:
            for series_item in data:
                series_name = series_item.get('name', '')
                x_vals = series_item.get('x', [])
                y_vals = series_item.get('y', [])
                size_vals = series_item.get('size', [])
                chart_data.add_series(series_name, list(zip(x_vals, y_vals, size_vals)))
                y_field_names.append(series_name)
        else:
            x_field = element.get('x', '_category')
            y_fields = element.get('y', [])
            if isinstance(y_fields, str):
                y_fields = [y_fields]
            names = element.get('names', y_fields)
            y_field_names = list(y_fields)
            for i, y_field in enumerate(y_fields):
                series_name = names[i] if i < len(names) else y_field
                if isinstance(size_field, list):
                    s_field = size_field[i] if i < len(size_field) else None
                else:
                    s_field = size_field
                points = []
                for row in data:
                    x_val = row.get(x_field, 0)
                    y_val = row.get(y_field, 0)
                    s_val = row.get(s_field, 0) if s_field else 0
                    points.append((float(x_val), float(y_val), float(s_val)))
                chart_data.add_series(series_name, points)
    else:
        chart_data = ChartData()
        if _is_series_oriented:
            categories = [str(v) for v in _first.get('x', [])]
            chart_data.categories = categories
            for series_item in data:
                series_name = series_item.get('name', '')
                values = series_item.get('y', [])
                chart_data.add_series(series_name, values)
                y_field_names.append(series_name)
        else:
            x_field = element.get('x', '_category')
            y_fields = element.get('y', [])
            if isinstance(y_fields, str):
                y_fields = [y_fields]
            names = element.get('names', y_fields)
            y_field_names = list(y_fields)
            chart_data.categories = [str(row.get(x_field, '')) for row in data] if x_field else []
            for i, y_field in enumerate(y_fields):
                series_name = names[i] if i < len(names) else y_field
                values = [row.get(y_field, 0) for row in data]
                chart_data.add_series(series_name, values)

    # --- Create chart shape (OPC skeleton only) ---
    chart_shape = slide.shapes.add_chart(
        xl_type, px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h), chart_data)

    # --- Build series data for XML builder ---
    series_data = []
    if _is_series_oriented:
        for i, series_item in enumerate(data):
            name = series_item.get('name', '')
            x_vals = [str(v) for v in series_item.get('x', [])]
            y_vals = series_item.get('y', [])
            series_data.append((name, x_vals, y_vals))
    else:
        x_field = element.get('x', '_category')
        categories_list = [str(row.get(x_field, '')) for row in data] if data else []
        for i, y_field in enumerate(y_field_names):
            name = y_field
            y_vals = [row.get(y_field, 0) for row in data]
            series_data.append((name, categories_list, y_vals))

    # --- Build complete chart XML via lxml ---
    chart_xml_bytes = _build_chart_xml(
        chart_type=chart_type,
        categories=[str(row.get(element.get('x', '_category'), '')) for row in data] if not _is_series_oriented else [str(v) for v in _first.get('x', [])],
        y_field_names=y_field_names,
        series_data=series_data,
        options=options,
        colors=element.get('colors', []),
        series_style=element.get('seriesStyle', {}),
        data_labels=element.get('dataLabels'),
        legend_spec=element.get('legend', _SENTINEL),
        title_spec=element.get('title'),
        x_axis=element.get('xAxis', {}),
        y_axis=element.get('yAxis', {}),
        secondary_axis=element.get('secondaryAxis'),
        theme=theme,
        font_family=options.get('fontFamily'),
    )

    # --- Overwrite chart part XML (replace _element to preserve OPC references) ---
    chart_part = chart_shape.chart.part
    chart_part._element = parse_xml(chart_xml_bytes)

    # --- Container styles ---
    fill_spec = element.get('fill')
    if fill_spec:
        apply_fill(chart_shape, fill_spec, theme, prs, base_dir)
    border_spec = element.get('border')
    if border_spec:
        apply_border(chart_shape, border_spec, theme)
    opacity = element.get('opacity')
    if opacity is not None and opacity < 1:
        apply_opacity(chart_shape, opacity)
    apply_transform(chart_shape, element)

    return chart_shape


# ---------------------------------------------------------------------------
# Background rendering
# ---------------------------------------------------------------------------

def _set_slide_background(slide, background, theme, prs, base_dir):
    """Set slide background from Fill spec."""
    if not background or not isinstance(background, dict):
        return
    bg_type = background.get('type')
    if bg_type == 'solid':
        color = resolve_color(background.get('color'), theme)
        if color:
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(color)
    elif bg_type == 'gradient':
        # Inject gradient into slide background XML
        bg = slide.background
        fill = bg.fill
        fill.solid()  # ensure spPr exists
        spPr = slide._element.cSld.find(qn('p:bg'))
        if spPr is None:
            bg_elem = slide._element.cSld.find(qn('p:bg'))
            if bg_elem is None:
                bg_elem = etree.SubElement(slide._element.cSld, qn('p:bg'))
            bgPr = bg_elem.find(qn('p:bgPr'))
            if bgPr is None:
                bgPr = etree.SubElement(bg_elem, qn('p:bgPr'))
            _inject_gradient_fill_bg(bgPr, background, theme)


def _inject_gradient_fill_bg(bgPr, fill_spec, theme):
    """Inject gradient fill into background properties."""
    stops = fill_spec.get('stops', [])
    if not stops:
        return
    angle = fill_spec.get('angle', 0)

    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
        for child in bgPr.findall(qn(tag)):
            bgPr.remove(child)

    grad_fill = etree.SubElement(bgPr, qn('a:gradFill'))
    gs_lst = etree.SubElement(grad_fill, qn('a:gsLst'))
    for stop in stops:
        pos = int(stop.get('position', 0) * 100000)
        color_hex = resolve_color(stop.get('color'), theme)
        if not color_hex:
            continue
        gs = etree.SubElement(gs_lst, qn('a:gs'))
        gs.set('pos', str(pos))
        color_clean = color_hex.lstrip('#')
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        # Check for alpha in hex string (8-char hex = #RRGGBBAA)
        if len(color_clean) > 6:
            srgb.set('val', color_clean[:6])
            alpha_8bit = int(color_clean[6:8], 16)
            alpha_ooxml = round(alpha_8bit / 255 * 100000)
            alpha_el = etree.SubElement(srgb, qn('a:alpha'))
            alpha_el.set('val', str(alpha_ooxml))
        else:
            srgb.set('val', color_clean)
    lin = etree.SubElement(grad_fill, qn('a:lin'))
    lin.set('ang', str(int(angle * 60000)))
    lin.set('scaled', '1')


# ---------------------------------------------------------------------------
# Template cleanup
# ---------------------------------------------------------------------------

def _cleanup_template_charts(pptx_path):
    """Remove all chart-related files from a copied template PPTX.

    python-pptx creates charts with only chartN.xml + embedded Excel.
    Templates may have additional chart style/color files (colorsN.xml,
    styleN.xml) that would become dangling references after new charts
    are created, causing PowerPoint to prompt for repair.
    """
    import zipfile
    import tempfile
    from xml.etree import ElementTree as ET

    chart_dir = 'ppt/charts/'
    embed_dir = 'ppt/embeddings/'

    # Collect Excel files referenced by chart relationships
    chart_excel_files = set()

    with zipfile.ZipFile(str(pptx_path), 'r') as zin:
        # Find chart Excel files from chart .rels files
        for name in zin.namelist():
            if name.startswith(chart_dir) and name.endswith('.rels'):
                try:
                    rels_xml = zin.read(name)
                    root = ET.fromstring(rels_xml)
                    for rel in root:
                        target = rel.get('Target', '')
                        if target.startswith('../embeddings/') and target.endswith('.xlsx'):
                            embed_name = 'ppt/' + target.lstrip('../')
                            chart_excel_files.add(embed_name)
                except Exception:
                    pass

        # Build list of entries to keep
        entries_to_copy = []
        for item in zin.infolist():
            # Remove entire ppt/charts/ directory
            if item.filename.startswith(chart_dir):
                continue
            # Remove chart-referenced Excel embeddings
            if item.filename in chart_excel_files:
                continue
            entries_to_copy.append(item)

        # Rewrite zip without chart files
        tmp_path = str(pptx_path) + '.tmp'
        with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in entries_to_copy:
                zout.writestr(item, zin.read(item.filename))

    import os
    os.replace(tmp_path, str(pptx_path))


# ---------------------------------------------------------------------------
# Main converter
# ---------------------------------------------------------------------------

def convert_pptd_to_pptx(pptd_path, output_path):
    pptd_path = Path(pptd_path).resolve()
    output_path = Path(output_path).resolve()

    if not pptd_path.exists():
        raise FileNotFoundError(f'PPTD file not found: {pptd_path}')

    base_dir = pptd_path.parent

    print(f'Loading: {pptd_path}')
    pptd_data, err = load_yaml(pptd_path)
    if err:
        raise ValueError(f'Failed to load PPTD: {err}')

    # Try to use source template for round-trip fidelity
    source_template = pptd_data.get('sourceTemplate')
    template_path = base_dir / source_template if source_template else None

    if template_path and template_path.exists():
        import shutil
        print(f'Using template: {template_path}')
        shutil.copy2(str(template_path), str(output_path))
        prs = Presentation(str(output_path))
        # Remove all existing slides from template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get(qn('r:id'))
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    else:
        prs = Presentation()

    # Use EMU sizes if available (exact round-trip), else convert from px
    slide_width = pptd_data.get('slideWidth')
    slide_height = pptd_data.get('slideHeight')
    if slide_width and slide_height:
        prs.slide_width = slide_width
        prs.slide_height = slide_height
    else:
        size = pptd_data.get('size', [1280, 720])
        prs.slide_width = px_to_emu(size[0])
        prs.slide_height = px_to_emu(size[1])

    theme = pptd_data.get('theme', {})

    # Resolve theme colors
    colors = theme.get('colors', {})
    resolved_colors = {}
    for key, val in colors.items():
        if isinstance(val, str) and val.startswith('$'):
            resolved_colors[key] = val
        else:
            resolved_colors[key] = val
    theme['colors'] = resolved_colors

    # Get blank layout (fallback)
    blank_layout = None
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower():
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    def _get_slide_layout(page_data):
        """Get the best slide layout for a page, using layoutIndex if available."""
        layout_idx = page_data.get('layoutIndex')
        if layout_idx is not None and 0 <= layout_idx < len(prs.slide_layouts):
            return prs.slide_layouts[layout_idx]
        return blank_layout

    pages = pptd_data.get('pages', [])
    total = len(pages)
    renderers = {
        'text': add_text_element,
        'shape': add_shape_element,
        'image': add_image_element,
        'icon': add_icon_element,
        'table': add_table_element,
        'chart': add_chart_element,
    }

    for page_idx, page_ref in enumerate(pages):
        page_path = base_dir / page_ref
        if not page_path.exists():
            print(f'  Warning: Page file not found: {page_path}')
            continue

        page_data, err = load_yaml(page_path)
        if err:
            print(f'  Warning: Failed to load {page_path}: {err}')
            continue

        slide = prs.slides.add_slide(_get_slide_layout(page_data))

        # Background
        background = page_data.get('background')
        if background:
            _set_slide_background(slide, background, theme, prs, base_dir)

        # Elements
        elements = page_data.get('elements', [])
        counts = {k: 0 for k in renderers}
        counts['skipped'] = 0

        # Track which placeholder indices are matched by PPTD elements
        matched_ph_indices = set()

        for element in elements:
            elem_type = element.get('elementType', 'text')
            renderer = renderers.get(elem_type)
            if renderer:
                try:
                    renderer(slide, element, theme, prs, base_dir)
                    counts[elem_type] += 1
                except Exception as e:
                    print(f'    Warning: Failed to render {element.get("elementId", "?")} ({elem_type}): {e}')
                    counts['skipped'] += 1
            else:
                counts['skipped'] += 1
            # Track matched placeholders
            ph_meta = element.get('placeholder')
            if ph_meta and ph_meta.get('idx') is not None:
                matched_ph_indices.add(ph_meta['idx'])

        # Remove unmatched layout placeholders (default content from template)
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                if ph_idx not in matched_ph_indices:
                    sp = shape._element
                    sp.getparent().remove(sp)

        # Notes
        notes = page_data.get('notes')
        if notes:
            try:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

        rendered = sum(v for k, v in counts.items() if k != 'skipped')
        print(f'  Page {page_idx + 1}/{total}: {page_ref} -> {rendered} elements '
              f"(text:{counts['text']}, shape:{counts['shape']}, image:{counts['image']}, "
              f"table:{counts['table']}, chart:{counts['chart']}, icon:{counts['icon']}, "
              f"skipped:{counts['skipped']})")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f'\nConversion complete!')
    print(f'  Output: {output_path}')
    print(f'  Total slides: {total}')
    final_size = pptd_data.get('size', [1280, 720])
    print(f'  Slide size: {final_size[0]}x{final_size[1]}px')

    return output_path


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser(description='Convert PPTD to PPTX format')
    parser.add_argument('input', help='Input .pptd file')
    parser.add_argument('-o', '--output', help='Output .pptx file (default: input.pptx)')
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f'Error: Input file not found: {args.input}')
        sys.exit(1)

    output = args.output or str(Path(args.input).with_suffix('.pptx'))
    try:
        convert_pptd_to_pptx(args.input, output)
    except Exception as e:
        print(f'Error: {e}')
        import traceback
        traceback.print_exc()
        sys.exit(1)
